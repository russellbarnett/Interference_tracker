"""
ELBOW ZONE™ | STRATEGIC BEHAVIORAL INVESTMENT TERMINAL
World-Class Presentation UI — Maximum Readability

CONFIDENTIAL: Russell Barnett © 2026. The Elbow Interference Theory™.
"""

# Suppress all warnings (Python 3.9 + Google SDK + urllib3)
import warnings
warnings.filterwarnings("ignore")

import streamlit as st
import pandas as pd
from dataclasses import dataclass
from typing import Optional, Tuple, Dict
import json

# ═══════════════════════════════════════════════════════════════════════════════
# INLINE TOOLTIP COPY (single-file, no extra modules)
# Streamlit shows this as a small "?" icon automatically on widgets that support `help=`
# ═══════════════════════════════════════════════════════════════════════════════

TOOLTIPS = {
    # Global
    "equation": "S = Satisfaction. Value Delivered ÷ Cost Extracted = Satisfaction. Satisfaction enables persistence over time, but persistence is not the equation itself.",
    "brand_name": "Enter the brand as a consumer recognizes it (typos are fine).",
    "format": "Format sets baseline friction. Portion-bound ends the occasion. Self-managed requires stopping decisions.",
    "occasion_free_text": "Optional. Adds context to the memo. Does not change math unless you use Priors.",

    # Six variables
    "M": "Mouthfeel: texture and physical experience while consuming.",
    "E": "Emotion: the payoff during consumption (comfort, indulgence, relief).",
    "F": "Familiarity: how established the brand is in real life (habit, ritual, legacy).",
    "B": "Bites: how many stop/continue decisions happen before the occasion ends.",
    "K": "Kinetic: physical effort to keep consuming (prep, tools, cleanup, handling).",
    "C": "Cognitive: how much thinking, evaluation, or self-control the product triggers.",

    # Overrides
    "override_rationale": "If you override a score, explain why. Keeps assumptions explicit and audit-able.",

    # Priors
    "priors_toggle": "Applies behavioral priors (context adjustments). Turn off for sensitivity checks.",
    "cohort": "Who is the buyer? Younger cohorts often carry higher messaging friction in values-heavy categories.",
    "priors_occasion": "Primary usage context. Occasion changes how friction and familiarity land.",
    "macro_stress": "When stress is active, Familiarity gets more weight (people default to known choices).",
    "promo_frequency": "Share of weeks on deal. High reliance implies velocity is being purchased.",
    "promo_depth": "Average discount depth. Deeper discount implies more purchased velocity risk.",
    "ups_13": "Units per store per week (last 13 weeks). Used for Repeat Momentum signal.",
    "ups_26": "Units per store per week (last 26 weeks). Used for Repeat Momentum signal.",

    # Upload
    "upload": "Upload CSV/Excel/PDF/Word/PPT/Text. Used for reference, context, and (optionally) priors inputs.",
}

def tip(key: str) -> str:
    return TOOLTIPS.get(key, "")

# Import priors module for behavioral adjustments
try:
    from priors import (
        apply_priors_dict, ScoringContext, get_default_context,
        calculate_repeat_momentum, generate_assumptions_section,
        RawScores, AdjustedScores
    )
    PRIORS_AVAILABLE = True
except ImportError:
    PRIORS_AVAILABLE = False
    print("[PRIORS] Module not available - using raw scores only")

# Import hard_data module for deterministic constants and guardrails
try:
    from hard_data import get_hard_data, get_external_support_snippets
    HARD_DATA = get_hard_data()
    HARD_DATA_AVAILABLE = True
    print("[HARD_DATA] Module loaded successfully")
except ImportError:
    HARD_DATA = None
    HARD_DATA_AVAILABLE = False
    print("[HARD_DATA] Module not available")

# ═══════════════════════════════════════════════════════════════════════════════════════
# GEMINI AI INTEGRATION
# ═══════════════════════════════════════════════════════════════════════════════════════

try:
    import google.generativeai as genai
    GEMINI_AVAILABLE = True
except ImportError:
    GEMINI_AVAILABLE = False

def get_gemini_model():
    """Initialize Gemini model with API key from secrets."""
    if not GEMINI_AVAILABLE:
        print("[GEMINI] Library not available")
        return None
    try:
        api_key = st.secrets.get("GEMINI_API_KEY", "")
        if not api_key or api_key == "YOUR_KEY_HERE":
            print("[GEMINI] No API key configured")
            return None
        genai.configure(api_key=api_key)
        model = genai.GenerativeModel('gemini-2.0-flash')
        print(f"[GEMINI] ✓ Model initialized (gemini-2.0-flash)")
        return model
    except Exception as e:
        print(f"[GEMINI] ✗ Error: {e}")
        return None

# The complete White Paper content for AI prompts
WHITE_PAPER_CONTEXT = """
THE ELBOW INTERFERENCE THEORY™ - FULL WHITE PAPER BY RUSSELL BARNETT

CORE CONCEPT:
Most products fail not because they lack quality, but because they ask too much of the consumer at the wrong moment.
Consumption begins in the elbow - the body acting without deliberation. Reaching, opening, tasting. The head is not involved.
But at some point, the head arrives. The consumer becomes aware they are consuming. They evaluate. They decide whether to continue.
The transition from elbow to head is the interference point. The longer the elbow runs uninterrupted, the more consumption compounds quietly.

KEY INSIGHT: Products do not compete on taste alone. They compete on how long they delay the head.

THE SATISFACTION EQUATION:
S = (M × E × F) ÷ (B × K × C)

VALUE DELIVERED (Numerator):
- M (Mouthfeel): How the product feels in the mouth during consumption
- E (Emotion): The feeling the product creates while being consumed  
- F (Familiarity): How recognizable and comfortable the experience is

COST EXTRACTED (Denominator):
- B (Bites): How many decisions the consumer makes before the occasion ends
- K (Kinetic Effort): The physical work required to continue consuming
- C (Cognitive Interference): How much the product asks the consumer to think

INTERPRETATION: S = Satisfaction. The option with higher S (Satisfaction) repeats more easily. The gap between scores tells you how large the behavioral advantage is.

CRITICAL INSIGHT ON FRICTION:
Friction compounds. Pleasure does not. A product rarely fails because of one flaw. It fails because multiple small frictions stack until the head interrupts the elbow.

CASE STUDY - MY/MOCHI (Structural Compounder):
- Handheld, portion-bound, no spoon, no stick, one motion from freezer to mouth
- Did NOT replace pints - served a different job (frozen snacking vs ice cream eating)
- Scores: M=4 (differentiated chew), E=4 (novel/surprising), F=3 (moderate initially)
- Denominator: B=1 (portion-bound), K=1 (single motion), C=1 (occasion ends automatically)
- S-Score: 16.0 - DENOMINATOR COLLAPSE achieved
- Key: "My/Mochi delays the head by ending the event"

CASE STUDY - PINT ICE CREAM (High Interference):
- Multiple bites, self-managed stopping, spoon/bowl ritual
- Scores: M=5 (category benchmark), E=5 (high indulgence), F=5 (very familiar)
- Denominator: B=4-5 (many decision points), K=3 (spoon/thaw/cleanup), C=3-4 (guilt, stopping decisions)
- S-Score: ~1.6-2.2
- Key: "The head arrives when stopping becomes self-managed"

THE CRITICAL DISTINCTION - FAMILIARITY (F):
- Legacy brands with decades of presence = F=5 (e.g., Serendipity has NYC ritual since 1954)
- New entrants/startups = F=2-3 (consumers must learn the product)
- Celebrity launches = F=3 (novelty, not ritual - the celebrity is familiar but the PRODUCT is not)

THE CRITICAL DISTINCTION - COGNITIVE INTERFERENCE (C):
- Legacy brands you buy on autopilot = C=1-2
- Premium pricing requiring justification = C=3-4  
- Celebrity brands = C=4-5 because "celebrity branding invites the Head Zone to audit the purchase"
- The Snoop Dogg hook is CONCEPTUAL (requires thinking). My/Mochi's success was EMBODIED (resolved in the mouth)
- "Conceptual hooks are harder and more expensive to defend over time"

THE "YEAR 5 SWING" RISK:
For high-interference brands, once novelty fades, the Head Zone must manually "authorize" the purchase. This creates a "Behavioral Margin Call" where velocity must be purchased through sustained spend rather than quiet compounding.

PRESENCE VS PERSISTENCE:
- TAM (Total Addressable Market) measures PRESENCE - how many consumers might consider a product
- S-Score = Satisfaction. Higher Satisfaction supports persistence (whether they will return without being reminded)
- "A $10 billion TAM means nothing if the product invites the head in too early"

ENTERPRISE VALUE INSIGHT:
"Products do not create behavior. They align with it. Capital can extend momentum, but it cannot change structure. Enterprise value compounds where the elbow finishes before the head arrives."
"""


def normalize_brand_name(name: str) -> str:
    """Normalize brand name for better AI recognition."""
    if not name:
        return name
    
    # Common variations mapping
    variations = {
        "dr bombay": "Dr. Bombay",
        "drbombay": "Dr. Bombay", 
        "dr.bombay": "Dr. Bombay",
        "doctor bombay": "Dr. Bombay",
        "mymochi": "My/Mochi",
        "my mochi": "My/Mochi",
        "cocacola": "Coca-Cola",
        "coca cola": "Coca-Cola",
        "coke": "Coca-Cola",
        "pepsicola": "Pepsi",
        "pepsi cola": "Pepsi",
        "lays": "Lay's",
        "lay's": "Lay's",
        "haagen dazs": "Häagen-Dazs",
        "haagendazs": "Häagen-Dazs",
        "haagen-dazs": "Häagen-Dazs",
        "ben and jerrys": "Ben & Jerry's",
        "ben & jerry's": "Ben & Jerry's",
        "ben jerrys": "Ben & Jerry's",
    }
    
    normalized = name.strip().lower()
    if normalized in variations:
        return variations[normalized]
    
    # Title case if not in variations
    return name.strip()


def analyze_brand_with_ai(brand_name: str) -> Optional[dict]:
    """
    Use Gemini to analyze a brand using the complete Elbow Interference Theory™.
    Returns dict with archetype, description, and BRAND-SPECIFIC scores.
    """
    model = get_gemini_model()
    if not model or not brand_name:
        return None
    
    # Normalize the brand name
    normalized_name = normalize_brand_name(brand_name)
    
    prompt = f"""{WHITE_PAPER_CONTEXT}

=== YOUR TASK ===

Analyze this CPG GROCERY brand using the Elbow Interference Theory: "{normalized_name}"

NOTE: The user entered "{brand_name}" - recognize common spelling variations:
- "dr bombay", "Dr. Bombay", "Dr Bombay" = Snoop Dogg's ice cream brand
- "mymochi", "My/Mochi", "my mochi" = the mochi ice cream brand
- Treat typos and variations as the intended brand

This applies to ANY orally consumed CPG grocery item - anything you eat or drink:
- Snacks (chips, crackers, cookies, bars, nuts, popcorn, jerky)
- Frozen (ice cream, novelties, pizza, meals, appetizers)
- Dairy (yogurt, cheese, milk, butter, cream)
- Beverages (soda, water, juice, coffee, tea, energy drinks, alcohol, kombucha)
- Confectionery (candy, chocolate, gum, mints)
- Bakery (bread, pastries, muffins, cakes)
- Deli/Prepared (sandwiches, salads, ready meals)
- Breakfast (cereal, oatmeal, pancake mix, syrup)
- Condiments (sauces, dressings, spreads)
- Baby food, supplements, protein powders - anything consumed orally

FOCUS ON F (Familiarity) and C (Cognitive) - these are BRAND-SPECIFIC:

F (Familiarity) - "How established is this brand in the consumer's mind?"

CRITICAL EXAMPLES FROM THE WHITE PAPER:
- **Serendipity** = F=5 (iconic NYC brand since 1954, decades of cultural presence, restaurant institution)
- **Dr. Bombay** = F=3 (launched 2023, celebrity-driven, consumers don't have RITUAL with it yet)
- **Häagen-Dazs** = F=5 (legacy premium brand, decades of presence)
- **My/Mochi** = F=3 (novel format, growing but not yet ritual)

General guidance:
- F=5: Iconic legacy brands - household names for DECADES (Coca-Cola, Oreo, Lay's, Serendipity, Ben & Jerry's)
- F=4: Well-known established brands (Kind, Chobani, Talenti)
- F=3: Growing/emerging brands OR celebrity launches (Poppi, Dr. Bombay, Feastables)
- F=2: New market entrants, niche brands
- F=1: Unknown/brand new launches

C (Cognitive) - "Does this brand make consumers THINK before buying?"

CRITICAL EXAMPLES FROM THE WHITE PAPER:
- **Serendipity** = C=3 (familiar legacy, some premium consideration but mostly autopilot)
- **Dr. Bombay** = C=5 (celebrity noise FORCES evaluation - "is this worth it because of Snoop?")
- **Coca-Cola** = C=1 (pure autopilot, decades of habit)
- **Häagen-Dazs** = C=3 (premium but familiar)

THE CELEBRITY TRAP: Celebrity brands get C=4-5 because:
> "Celebrity branding invites the Head Zone to audit the purchase"
> The consumer evaluates the CELEBRITY, not just the food

General guidance:
- C=1-2: Autopilot purchases - grab without thinking (Coke, Oreo, legacy brands)
- C=3: Some consideration - premium positioning (Häagen-Dazs, Serendipity)
- C=4-5: High cognitive load - CELEBRITY BRANDS, unfamiliar concepts (Dr. Bombay, Feastables)

CELEBRITY BRANDS ARE SPECIAL:
- Celebrity association INCREASES cognitive load (C=4-5) because consumers evaluate the celebrity, not just the food
- Celebrity does NOT increase familiarity (F) - the celebrity is familiar, but the PRODUCT is new
- Examples: Dr. Bombay (Snoop Dogg), Feastables (MrBeast), Prime (Logan Paul) = F=2-3, C=4-5

ALSO PROVIDE M AND E (but these matter less - focus on F and C):
- M (Mouthfeel): Quality of sensory experience (1-5)
- E (Emotion): Satisfaction/indulgence delivered (1-5)

For B and K, estimate based on typical format (user will override with their format selection):

UNITIZED formats (B=1, K=1) - The package ends the occasion:
- Single-serve packs, bars, sticks, individual portions
- Single cans/bottles of beverages
- Individually wrapped items
- Single-serve cups, pouches, squeeze tubes

BULK formats (B=4, K=3) - Consumer decides when to stop:
- Bags, pouches (chips, crackers, cookies, candy)
- Tubs, pints, containers (ice cream, yogurt, dips)
- Boxes, cartons (cereal, crackers, snack mixes)
- Multi-use jars/bottles (peanut butter, sauces, condiments)
- Anything where you reach in multiple times

Return ONLY valid JSON:
{{"archetype": "unitized|bulk|ritual", "description": "brief description", "M": 4, "E": 4, "F": 3, "B": 4, "K": 3, "C": 4, "reasoning": "Explain your F and C scores - what makes this brand familiar or unfamiliar? What creates cognitive load?"}}

If you cannot identify the brand, return: {{"error": "unknown"}}"""

    try:
        print(f"\n{'='*60}")
        print(f"🔍 ANALYZING BRAND: {brand_name} (normalized: {normalized_name})")
        print(f"{'='*60}")
        
        response = model.generate_content(prompt)
        text = response.text.strip()
        
        print(f"\n📥 RAW AI RESPONSE:")
        print(f"{text}")
        print(f"\n{'='*60}")
        
        if "```json" in text:
            text = text.split("```json")[1].split("```")[0]
        elif "```" in text:
            text = text.split("```")[1].split("```")[0]
        
        result = json.loads(text.strip())
        
        print(f"\n✅ PARSED SCORES FOR {brand_name}:")
        print(f"   M (Mouthfeel)  = {result.get('M')}")
        print(f"   E (Emotion)    = {result.get('E')}")
        print(f"   F (Familiarity)= {result.get('F')} ⭐")
        print(f"   B (Bites)      = {result.get('B')}")
        print(f"   K (Kinetic)    = {result.get('K')}")
        print(f"   C (Cognitive)  = {result.get('C')} ⭐")
        print(f"   Reasoning: {result.get('reasoning', 'none')[:100]}...")
        print(f"{'='*60}\n")
        
        return result
    except Exception as e:
        print(f"\n❌ AI ANALYSIS ERROR for {brand_name}: {e}")
        print(f"{'='*60}\n")
        return None

def generate_strategic_synthesis(brand1_name: str, brand2_name: str, 
                                  s1: float, s2: float,
                                  scores1: dict, scores2: dict,
                                  rationale1: str, rationale2: str) -> str:
    """
    Use Gemini to generate an institutional-grade investment memo.
    Produces comprehensive analysis like the Suu Tri-Biotics example.
    """
    print(f"[SYNTHESIS] Getting Gemini model...")
    model = get_gemini_model()
    if not model:
        print(f"[SYNTHESIS] ✗ Model is None!")
        return None
    print(f"[SYNTHESIS] ✓ Model acquired")
    
    num1 = scores1.get('M',1) * scores1.get('E',1) * scores1.get('F',1)
    den1 = scores1.get('B',1) * scores1.get('K',1) * scores1.get('C',1)
    num2 = scores2.get('M',1) * scores2.get('E',1) * scores2.get('F',1)
    den2 = scores2.get('B',1) * scores2.get('K',1) * scores2.get('C',1)
    
    # Determine which brand is the structural winner
    winner = brand1_name if s1 > s2 else brand2_name
    loser = brand2_name if s1 > s2 else brand1_name
    winner_score = max(s1, s2)
    loser_score = min(s1, s2)
    winner_scores = scores1 if s1 > s2 else scores2
    loser_scores = scores2 if s1 > s2 else scores1
    ratio = winner_score / loser_score if loser_score > 0 else float('inf')
    
    # Identify critical weaknesses (scores >= 4 in denominator)
    b1_warnings = []
    if scores1.get('B', 1) >= 4: b1_warnings.append(f"B={scores1.get('B')} (high decision count)")
    if scores1.get('K', 1) >= 4: b1_warnings.append(f"K={scores1.get('K')} (high effort)")
    if scores1.get('C', 1) >= 4: b1_warnings.append(f"C={scores1.get('C')} (head arrives early)")
    
    b2_warnings = []
    if scores2.get('B', 1) >= 4: b2_warnings.append(f"B={scores2.get('B')} (high decision count)")
    if scores2.get('K', 1) >= 4: b2_warnings.append(f"K={scores2.get('K')} (high effort)")
    if scores2.get('C', 1) >= 4: b2_warnings.append(f"C={scores2.get('C')} (head arrives early)")
    
    prompt = f"""{WHITE_PAPER_CONTEXT}

=== YOUR TASK ===

Generate a COMPREHENSIVE ELBOW INTERFERENCE ANALYSIS for Matt Leeds at Forward Consumer Partners.

This should be institutional-grade analysis that breaks down EACH brand individually, similar to how an analyst would deconstruct a brand like Suu Tri-Biotics against the framework.

=== BRAND DATA ===

**{brand1_name}**
- S-Score™: {s1:.2f}
- M={scores1.get('M',0)}, E={scores1.get('E',0)}, F={scores1.get('F',0)} (Numerator: {num1})
- B={scores1.get('B',0)}, K={scores1.get('K',0)}, C={scores1.get('C',0)} (Denominator: {den1})
- Format/Archetype: {scores1.get('archetype', 'unknown')}
{f"- ⚠️ Critical Weaknesses: {', '.join(b1_warnings)}" if b1_warnings else "- ✅ No critical denominator weaknesses"}
{f"- Analyst Note: {rationale1}" if rationale1 else ""}

**{brand2_name}**
- S-Score™: {s2:.2f}
- M={scores2.get('M',0)}, E={scores2.get('E',0)}, F={scores2.get('F',0)} (Numerator: {num2})
- B={scores2.get('B',0)}, K={scores2.get('K',0)}, C={scores2.get('C',0)} (Denominator: {den2})
- Format/Archetype: {scores2.get('archetype', 'unknown')}
{f"- ⚠️ Critical Weaknesses: {', '.join(b2_warnings)}" if b2_warnings else "- ✅ No critical denominator weaknesses"}
{f"- Analyst Note: {rationale2}" if rationale2 else ""}

**Structural Winner: {winner}** with {ratio:.2f}x behavioral advantage

=== OUTPUT FORMAT (FOLLOW THIS EXACTLY) ===

---

# ELBOW INTERFERENCE ANALYSIS: {brand1_name} vs {brand2_name}

## Product Structure Overview

Create a brief table comparing the two brands:
| Dimension | {brand1_name} | {brand2_name} |
|-----------|--------------|---------------|
| Format | ? | ? |
| Category | ? | ? |
| Consumption | ? | ? |

---

## {brand1_name}: DETAILED ANALYSIS

### Value Delivered (Numerator)

**M (Mouthfeel): {scores1.get('M',0)}/5**
- Explain what this score means for THIS specific brand
- Is it category-defining or merely acceptable?

**E (Emotion): {scores1.get('E',0)}/5**
- What emotional payoff does this brand provide?
- Is it immediate pleasure or delayed gratification?

**F (Familiarity): {scores1.get('F',0)}/5** {"⭐" if scores1.get('F',0) >= 4 else "⚠️" if scores1.get('F',0) <= 2 else ""}
- Is this anchored to existing ritual or requiring new behavior?
- Legacy vs new entrant assessment

### Cost Extracted (Denominator)

**B (Bites/Decisions): {scores1.get('B',0)}/5** {"⚠️ HIGH FRICTION" if scores1.get('B',0) >= 4 else ""}
- How many decisions before occasion ends?
- Does occasion end automatically or require self-management?

**K (Kinetic Effort): {scores1.get('K',0)}/5** {"⚠️ HIGH FRICTION" if scores1.get('K',0) >= 4 else ""}
- Physical effort to continue consuming
- Format-driven assessment

**C (Cognitive Interference): {scores1.get('C',0)}/5** {"⚠️ CRITICAL WEAKNESS" if scores1.get('C',0) >= 4 else ""}
- When does the head arrive?
- Celebrity noise? Premium justification? Conceptual hooks?

---

## {brand2_name}: DETAILED ANALYSIS

(Same structure as above for the second brand)

---

## The Core Problem: Head vs Elbow

Using the White Paper framework, identify the critical structural difference between these brands. Quote the relevant passage:

> "Consumption begins in the elbow. The elbow is the body acting without deliberation..."

For the LOWER scoring brand ({loser}), explain:
- Does it pass the "automatic consumption" test?
- At what point does the head arrive?
- Is the brand hook Conceptual (requires thinking) or Embodied (resolved in the body)?

---

## Comparative Scorecard

| Dimension | {brand1_name} | {brand2_name} | Advantage |
|-----------|--------------|---------------|-----------|
| Ritual Preservation | Yes/No | Yes/No | ? |
| Elbow Entry | How? | How? | ? |
| Cognitive Load | Low/Med/High | Low/Med/High | ? |
| Head Arrival | When? | When? | ? |
| Repeat Mechanism | How? | How? | ? |

Key Quote: *"Products do not compete on taste alone. They compete on how long they delay the head."*

---

## Why This Matters for Enterprise Value

The White Paper states: *"TAM is the size of the door. Elbow Interference tells you whether consumers walk through it once or every week."*

For {loser} (S-Score: {loser_score:.2f}):
- High interference = velocity must be purchased, not compounded
- Repeat depends on ______, not automatic behavior
- Marketing must continuously reactivate
- Churn risk assessment

For {winner} (S-Score: {winner_score:.2f}):
- Lower interference = repeat compounds naturally
- Behavioral advantage creates moat

Compare to benchmark acquisitions (Poppi at $1.95B captured low-friction soda moment).

---

## Where {loser} Could Improve (Through EIT Lens)

Provide 3-5 specific, actionable recommendations using the framework:
1. How to lower B (decisions)?
2. How to lower K (effort)?
3. How to lower C (cognitive friction)?
4. How to increase F (familiarity)?
5. What format change would collapse the denominator?

---

## Final Assessment

Using the framework's language:

> "Repeat behavior is not created by marketing. It is allowed or blocked by product structure."

**{winner}**: [Structural Compounder / Low-Friction Asset]
- Why repeat is structurally enabled

**{loser}**: [High-Interference Asset / Growth Trap / Purchased Velocity]
- Why repeat requires continued investment

**Investment Verdict:**
- Clear recommendation for Matt Leeds
- S-Score ratio ({ratio:.2f}x) interpretation
- What would change the calculus?

---

Write in institutional investment tone. Be specific to these brands. Use actual scores throughout. Include White Paper quotes where impactful."""

    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        print(f"Memo generation error: {e}")
        return None


def generate_rule_based_memo(brand1: str, brand2: str, s1: float, s2: float, 
                              scores1: dict, scores2: dict) -> str:
    """
    Generate a comprehensive rule-based analysis memo without AI.
    Uses the White Paper logic directly - styled like the Suu Tri-Biotics example.
    """
    winner = brand1 if s1 > s2 else brand2
    loser = brand2 if s1 > s2 else brand1
    winner_score = max(s1, s2)
    loser_score = min(s1, s2)
    winner_scores = scores1 if s1 > s2 else scores2
    loser_scores = scores2 if s1 > s2 else scores1
    ratio = winner_score / loser_score if loser_score > 0 else float('inf')
    
    num1 = scores1['M'] * scores1['E'] * scores1['F']
    den1 = scores1['B'] * scores1['K'] * scores1['C']
    num2 = scores2['M'] * scores2['E'] * scores2['F']
    den2 = scores2['B'] * scores2['K'] * scores2['C']
    
    # Determine key differentiators
    f_diff = scores1['F'] - scores2['F']
    c_diff = scores2['C'] - scores1['C']
    
    # Warning flags
    def get_warnings(scores, name):
        warnings = []
        if scores['B'] >= 4: warnings.append(f"B={scores['B']} ⚠️ HIGH (many decisions)")
        if scores['K'] >= 4: warnings.append(f"K={scores['K']} ⚠️ HIGH (effort required)")
        if scores['C'] >= 4: warnings.append(f"C={scores['C']} ⚠️ CRITICAL (head arrives early)")
        return warnings
    
    b1_warnings = get_warnings(scores1, brand1)
    b2_warnings = get_warnings(scores2, brand2)
    
    # Build comprehensive memo
    memo = f"""# ELBOW INTERFERENCE ANALYSIS: {brand1} vs {brand2}

---

## Executive Summary

**{winner}** holds a **{ratio:.1f}x structural advantage** over {loser} in behavioral persistence.

Using the Satisfaction Equation: **S = (M × E × F) ÷ (B × K × C)**

| Brand | S-Score™ | Numerator | Denominator |
|-------|----------|-----------|-------------|
| {brand1} | **{s1:.2f}** | {num1} | {den1} |
| {brand2} | **{s2:.2f}** | {num2} | {den2} |

---

## {brand1}: Detailed Analysis

### Value Delivered (Numerator = {num1})

**M (Mouthfeel): {scores1['M']}/5**
- Sensory experience during consumption
- {"Strong sensory delivery" if scores1['M'] >= 4 else "Moderate sensory experience" if scores1['M'] >= 3 else "Limited mouthfeel impact"}

**E (Emotion): {scores1['E']}/5**
- Emotional satisfaction during occasion
- {"High emotional payoff - immediate pleasure" if scores1['E'] >= 4 else "Moderate emotional connection" if scores1['E'] >= 3 else "Low emotional engagement"}

**F (Familiarity): {scores1['F']}/5** {"⭐ STRENGTH" if scores1['F'] >= 4 else "⚠️ RISK" if scores1['F'] <= 2 else ""}
- {"Anchored to established ritual - autopilot purchasing" if scores1['F'] >= 4 else "Some brand recognition" if scores1['F'] >= 3 else "Requires building new behavior - unfamiliar to consumers"}

### Cost Extracted (Denominator = {den1})

**B (Bites/Decisions): {scores1['B']}/5** {"⚠️ HIGH FRICTION" if scores1['B'] >= 4 else ""}
- {"Many decisions before occasion ends - requires self-management" if scores1['B'] >= 4 else "Moderate decision count" if scores1['B'] >= 2 else "Occasion ends automatically - minimal decisions"}

**K (Kinetic Effort): {scores1['K']}/5** {"⚠️ HIGH FRICTION" if scores1['K'] >= 4 else ""}
- {"Significant physical effort required" if scores1['K'] >= 4 else "Moderate effort" if scores1['K'] >= 2 else "Minimal effort - grab and consume"}

**C (Cognitive Interference): {scores1['C']}/5** {"⚠️ CRITICAL WEAKNESS" if scores1['C'] >= 4 else ""}
- {"Head arrives before elbow finishes - consumer must think/evaluate" if scores1['C'] >= 4 else "Some cognitive processing required" if scores1['C'] >= 3 else "Autopilot consumption - elbow runs uninterrupted"}

{f"**Denominator Warnings:** {', '.join(b1_warnings)}" if b1_warnings else "✅ **No critical denominator weaknesses**"}

---

## {brand2}: Detailed Analysis

### Value Delivered (Numerator = {num2})

**M (Mouthfeel): {scores2['M']}/5**
- Sensory experience during consumption
- {"Strong sensory delivery" if scores2['M'] >= 4 else "Moderate sensory experience" if scores2['M'] >= 3 else "Limited mouthfeel impact"}

**E (Emotion): {scores2['E']}/5**
- Emotional satisfaction during occasion
- {"High emotional payoff - immediate pleasure" if scores2['E'] >= 4 else "Moderate emotional connection" if scores2['E'] >= 3 else "Low emotional engagement"}

**F (Familiarity): {scores2['F']}/5** {"⭐ STRENGTH" if scores2['F'] >= 4 else "⚠️ RISK" if scores2['F'] <= 2 else ""}
- {"Anchored to established ritual - autopilot purchasing" if scores2['F'] >= 4 else "Some brand recognition" if scores2['F'] >= 3 else "Requires building new behavior - unfamiliar to consumers"}

### Cost Extracted (Denominator = {den2})

**B (Bites/Decisions): {scores2['B']}/5** {"⚠️ HIGH FRICTION" if scores2['B'] >= 4 else ""}
- {"Many decisions before occasion ends - requires self-management" if scores2['B'] >= 4 else "Moderate decision count" if scores2['B'] >= 2 else "Occasion ends automatically - minimal decisions"}

**K (Kinetic Effort): {scores2['K']}/5** {"⚠️ HIGH FRICTION" if scores2['K'] >= 4 else ""}
- {"Significant physical effort required" if scores2['K'] >= 4 else "Moderate effort" if scores2['K'] >= 2 else "Minimal effort - grab and consume"}

**C (Cognitive Interference): {scores2['C']}/5** {"⚠️ CRITICAL WEAKNESS" if scores2['C'] >= 4 else ""}
- {"Head arrives before elbow finishes - consumer must think/evaluate" if scores2['C'] >= 4 else "Some cognitive processing required" if scores2['C'] >= 3 else "Autopilot consumption - elbow runs uninterrupted"}

{f"**Denominator Warnings:** {', '.join(b2_warnings)}" if b2_warnings else "✅ **No critical denominator weaknesses**"}

---

## The Core Problem: Head vs Elbow

> *"Consumption begins in the elbow. The elbow is the body acting without deliberation. Reaching, opening, tasting. The head is not involved."*

"""
    
    # Determine which brand has the head problem
    if loser_scores['C'] >= 4:
        memo += f"""**{loser} invites the Head immediately.**

With C={loser_scores['C']}, the consumer must THINK before/during consumption. This creates:
- Continuous self-interrogation ("Is this worth it?")
- Decision fatigue on repeat purchase
- Vulnerability to competitive switching

The White Paper is clear: *"The moment the head arrives, the elbow slows."*
"""
    elif max(loser_scores['B'], loser_scores['K']) >= 4:
        memo += f"""**{loser} creates friction through effort.**

With B={loser_scores['B']} and K={loser_scores['K']}, the physical structure of consumption creates repeated decision points. Each decision is an opportunity for the head to arrive.
"""
    else:
        memo += f"""Both brands have manageable interference profiles. The difference lies primarily in the numerator (value delivered) rather than denominator (cost extracted).
"""
    
    # Comparative table
    memo += f"""

---

## Comparative Scorecard

| Dimension | {brand1} | {brand2} |
|-----------|----------|----------|
| Elbow Entry | {"Easy - minimal prep" if scores1['K'] <= 2 else "Moderate" if scores1['K'] <= 3 else "Requires effort"} | {"Easy - minimal prep" if scores2['K'] <= 2 else "Moderate" if scores2['K'] <= 3 else "Requires effort"} |
| Cognitive Load | {"Low" if scores1['C'] <= 2 else "Medium" if scores1['C'] <= 3 else "High"} | {"Low" if scores2['C'] <= 2 else "Medium" if scores2['C'] <= 3 else "High"} |
| Head Arrival | {"Delayed" if scores1['C'] <= 2 else "Mid-occasion" if scores1['C'] <= 3 else "Early/Immediate"} | {"Delayed" if scores2['C'] <= 2 else "Mid-occasion" if scores2['C'] <= 3 else "Early/Immediate"} |
| Repeat Mechanism | {"Habitual" if den1 <= 6 else "Requires effort"} | {"Habitual" if den2 <= 6 else "Requires effort"} |

---

## Why This Matters for Enterprise Value

> *"TAM is the size of the door. Elbow Interference tells you whether consumers walk through it once or every week."*

"""
    
    # Add specific enterprise value analysis
    if loser_score < 2:
        memo += f"""**{loser}** (S-Score: {loser_score:.2f}):
- High interference = velocity must be **purchased**, not compounded
- Repeat depends on consumer **discipline**, not automatic behavior
- Marketing must continuously **reactivate** ("Have you bought yours lately?")
- **Churn risk** when consumers forget, get busy, or encounter friction

**{winner}** (S-Score: {winner_score:.2f}):
- Lower interference = repeat **compounds naturally**
- Behavioral advantage creates **structural moat**
- Marketing drives **trial**; structure handles **repeat**
"""
    else:
        memo += f"""Both brands show viable persistence profiles, but **{winner}** has the structural edge.

- {winner}: Velocity earned through repeat structure
- {loser}: Requires more marketing support for same repeat rate
"""
    
    # Recommendations for the weaker brand
    memo += f"""

---

## Where {loser} Could Improve (Through EIT Lens)

"""
    if loser_scores['C'] >= 4:
        memo += f"""1. **Lower Cognitive Friction**: Simplify the value proposition. Stop requiring consumers to think.
2. **Reduce Conceptual Hooks**: If brand relies on celebrity, trend, or complex benefits, find ways to make value embodied (felt immediately) rather than conceptual (requires explanation).
"""
    if loser_scores['B'] >= 4:
        memo += f"""3. **Reduce Decisions**: Consider portion-controlled or single-serve formats that end the occasion automatically.
"""
    if loser_scores['K'] >= 3:
        memo += f"""4. **Lower Kinetic Effort**: Ready-to-consume formats beat "some assembly required."
"""
    if loser_scores['F'] <= 2:
        memo += f"""5. **Build Familiarity**: Anchor to existing rituals rather than creating new behaviors.
"""
    
    # Final assessment
    memo += f"""

---

## Final Assessment

> *"Repeat behavior is not created by marketing. It is allowed or blocked by product structure."*

**{winner}**: {"Structural Compounder" if winner_score >= 5 else "Low-Friction Asset" if winner_score >= 2 else "Moderate Persistence"}
- {"Denominator collapse achieved - repeat is automatic" if winner_scores['B'] * winner_scores['K'] * winner_scores['C'] <= 6 else "Manageable friction profile"}
- Velocity {"compounds naturally" if winner_score >= 3 else "supported by structure"}

**{loser}**: {"High-Interference Asset" if loser_score < 1.5 else "Purchased Velocity Risk" if loser_score < 2.5 else "Moderate Persistence"}
- {"Heavy denominator blocks habitual repeat" if loser_scores['B'] * loser_scores['K'] * loser_scores['C'] >= 24 else "Some structural friction"}
- {"Velocity must be purchased through continuous marketing" if loser_score < 2 else "Requires some marketing support"}

**Investment Verdict:**

{winner} wins this duel with a **{ratio:.1f}x structural advantage**. 

{"This is a decisive gap - the structural difference is significant enough to impact long-term enterprise value." if ratio >= 2 else "This is a meaningful gap that compounds over time through repeat purchase behavior." if ratio >= 1.5 else "The gap is modest - competitive positioning and execution will matter as much as structure."}

---

*Analysis generated using the Elbow Interference Theory™ by Russell Barnett*
"""
    
    return memo

# ═══════════════════════════════════════════════════════════════════════════════════════
# CONFIGURATION
# ═══════════════════════════════════════════════════════════════════════════════════════

st.set_page_config(
    page_title="Elbow Interference Evaluator™",
    page_icon="◈",
    layout="wide",
    initial_sidebar_state="expanded"
)

# ═══════════════════════════════════════════════════════════════════════════════════════
# PRODUCT FORMATS (Archetypes)
# ═══════════════════════════════════════════════════════════════════════════════════════
#
# The PHYSICAL FORMAT of a product determines how much "interference" occurs during
# consumption. The format sets baseline scores for the Satisfaction Equation.
#
# KEY INSIGHT: The Denominator (B × K × C) measures FRICTION.
#   - B (Bites) = How many decisions before you're done?
#   - K (Kinetic) = How much physical effort to keep eating?
#   - C (Cognitive) = How much thinking required?
#
# LOW Denominator = Product "does the work" = Easy to finish = High repeat
# HIGH Denominator = You must "manage yourself" = Hard to finish = Low repeat
# ═══════════════════════════════════════════════════════════════════════════════════════

@dataclass(frozen=True)
class Archetype:
    name: str           # Display name
    short_name: str     # Simple label
    examples: str       # Real product examples
    description: str    # Plain English explanation
    M: int              # Mouthfeel (1-5)
    E: int              # Emotion (1-5)
    F: int              # Familiarity (1-5)
    B: int              # Bites/Decisions (1-5)
    K: int              # Kinetic/Effort (1-5)
    C: int              # Cognitive/Thinking (1-5)
    why_denominator: str  # Why the denominator is set this way


ARCHETYPES: Dict[str, Archetype] = {
    "unitized": Archetype(
        name="Single-Serve / Portion-Controlled",
        short_name="Single-Serve",
        examples="Mochi bites, protein bars, ice cream bars, snack packs",
        description="The package decides when you're done. One portion = one occasion. No decisions needed.",
        M=4, E=4, F=3, B=1, K=1, C=1,
        why_denominator="B=1, K=1, C=1 because you eat the whole thing without deciding to stop. The format does the work."
    ),
    "bulk": Archetype(
        name="Multi-Serve / Open Container",
        short_name="Multi-Serve",
        examples="Pints, tubs, chip bags, family-size containers",
        description="YOU decide when to stop. Each bite is a new decision. More chances to quit early.",
        M=5, E=4, F=4, B=4, K=3, C=3,
        why_denominator="B=4, K=3, C=3 because you must self-regulate. Many bites = many chances for your brain to say 'stop'."
    ),
    "ritual": Archetype(
        name="Ritual / Single-Occasion Beverage",
        short_name="Ritual Drink",
        examples="Soda cans, energy drinks, bottled beverages, sparkling water",
        description="One can = one ritual. Highly familiar. You drink it without thinking.",
        M=3, E=3, F=5, B=1, K=1, C=1,
        why_denominator="B=1, K=1, C=1 because it's automatic. High familiarity (F=5) means zero cognitive friction."
    ),
}

BRAND_DATABASE: Dict[str, Tuple[str, str]] = {
    # Unitized brands (with common variations)
    "my/mochi": ("unitized", "Premium frozen mochi"), 
    "mymochi": ("unitized", "Premium frozen mochi"),
    "my mochi": ("unitized", "Premium frozen mochi"),
    "mochi": ("unitized", "Premium frozen mochi"),
    "kind": ("unitized", "Nutrition bars"), 
    "kind bar": ("unitized", "Nutrition bars"),
    "rxbar": ("unitized", "Protein bars"),
    "rx bar": ("unitized", "Protein bars"),
    "quest": ("unitized", "Protein snacks"), 
    "yasso": ("unitized", "Frozen yogurt bars"),
    "magnum": ("unitized", "Ice cream bars"), 
    "doughlicious": ("unitized", "Cookie dough bites"),
    "doughliscious": ("unitized", "Cookie dough bites"),  # Common misspelling
    
    # Bulk brands (with common variations)
    "ben & jerry's": ("bulk", "Premium pints"), 
    "ben and jerry's": ("bulk", "Premium pints"),
    "ben and jerrys": ("bulk", "Premium pints"),
    "ben jerry": ("bulk", "Premium pints"),
    "häagen-dazs": ("bulk", "Premium pints"),
    "haagen-dazs": ("bulk", "Premium pints"),
    "haagen dazs": ("bulk", "Premium pints"),
    "hagendaz": ("bulk", "Premium pints"),
    "haagendazs": ("bulk", "Premium pints"),
    "talenti": ("bulk", "Gelato pints"),
    "dr. bombay": ("bulk", "Celebrity pints"), 
    "dr bombay": ("bulk", "Celebrity pints"),
    "doctor bombay": ("bulk", "Celebrity pints"),
    "drbombay": ("bulk", "Celebrity pints"),
    "serendipity": ("bulk", "Premium pints"),
    "jeni's": ("bulk", "Artisan pints"), 
    "jenis": ("bulk", "Artisan pints"),
    "lay's": ("bulk", "Multi-serve chips"),
    "lays": ("bulk", "Multi-serve chips"),
    "doritos": ("bulk", "Multi-serve chips"),
    "cheetos": ("bulk", "Multi-serve snacks"),
    
    # Ritual brands (with common variations)
    "poppi": ("ritual", "Prebiotic soda"), 
    "poppie": ("ritual", "Prebiotic soda"),
    "olipop": ("ritual", "Functional soda"),
    "oli pop": ("ritual", "Functional soda"),
    "liquid death": ("ritual", "Canned water"), 
    "liquiddeath": ("ritual", "Canned water"),
    "celsius": ("ritual", "Energy drinks"),
    "red bull": ("ritual", "Energy drinks"),
    "redbull": ("ritual", "Energy drinks"),
    "coca-cola": ("ritual", "Carbonated beverages"),
    "coca cola": ("ritual", "Carbonated beverages"),
    "coke": ("ritual", "Carbonated beverages"),
    "pepsi": ("ritual", "Carbonated beverages"),
    "monster": ("ritual", "Energy drinks"),
}

# Single source of truth for known brand scores (M,E,F,B,K,C) and reasoning. Used by Brand 1 and Brand 2.
KNOWN_BRANDS = {
    'serendipity': {'M': 4, 'E': 4, 'F': 5, 'B': 4, 'K': 3, 'C': 3,
                    'reasoning': 'Serendipity: NYC legacy brand since 1954. High Familiarity (F=5) from decades of brand building. Standard pint format drives B/K. Low cognitive friction (C=3) - no justification needed for a known indulgence.'},
    'dr. bombay': {'M': 4, 'E': 4, 'F': 3, 'B': 4, 'K': 3, 'C': 5,
                   'reasoning': 'Dr. Bombay: Celebrity brand (Snoop Dogg). Low Familiarity (F=3) - new entrant without legacy. HIGH Cognitive (C=5) - celebrity branding invites consumer to evaluate/justify purchase. The "head" arrives before consumption.'},
    'dr bombay': {'M': 4, 'E': 4, 'F': 3, 'B': 4, 'K': 3, 'C': 5,
                  'reasoning': 'Dr. Bombay: Celebrity brand (Snoop Dogg). Low Familiarity (F=3) - new entrant without legacy. HIGH Cognitive (C=5) - celebrity branding invites consumer to evaluate/justify purchase.'},
    'häagen-dazs': {'M': 5, 'E': 5, 'F': 5, 'B': 4, 'K': 3, 'C': 2,
                    'reasoning': 'Häagen-Dazs: Ultra-legacy premium brand. Maximum Familiarity (F=5). Premium positioning but so established that C is low - consumers know what they are getting.'},
    'haagen-dazs': {'M': 5, 'E': 5, 'F': 5, 'B': 4, 'K': 3, 'C': 2,
                    'reasoning': 'Häagen-Dazs: Ultra-legacy premium brand.'},
    'ben & jerry\'s': {'M': 5, 'E': 5, 'F': 5, 'B': 4, 'K': 3, 'C': 2,
                       'reasoning': 'Ben & Jerry\'s: Iconic legacy brand with decades of equity.'},
    'my/mochi': {'M': 4, 'E': 4, 'F': 4, 'B': 1, 'K': 1, 'C': 1,
                 'reasoning': 'My/Mochi: Unitized format = Denominator Collapse. B=1, K=1, C=1 because occasion ends automatically.'},
    'mymochi': {'M': 4, 'E': 4, 'F': 4, 'B': 1, 'K': 1, 'C': 1,
                'reasoning': 'My/Mochi: Unitized format = Denominator Collapse.'},
    'coca-cola': {'M': 4, 'E': 4, 'F': 5, 'B': 1, 'K': 1, 'C': 1,
                  'reasoning': 'Coca-Cola: Maximum legacy (F=5). Ritual format (can) = Denominator Collapse.'},
    'coke': {'M': 4, 'E': 4, 'F': 5, 'B': 1, 'K': 1, 'C': 1,
             'reasoning': 'Coca-Cola: Maximum legacy.'},
    'pepsi': {'M': 4, 'E': 4, 'F': 5, 'B': 1, 'K': 1, 'C': 1,
              'reasoning': 'Pepsi: Maximum legacy (F=5). Ritual format.'},
    'oreo': {'M': 4, 'E': 4, 'F': 5, 'B': 2, 'K': 1, 'C': 1,
             'reasoning': 'Oreo: Iconic snack brand. Maximum Familiarity.'},
    'poppi': {'M': 3, 'E': 3, 'F': 4, 'B': 1, 'K': 1, 'C': 2,
              'reasoning': 'Poppi: Growing familiarity. Ritual preservation (soda format). Some cognitive for "better-for-you" positioning.'},
    'lays': {'M': 4, 'E': 4, 'F': 5, 'B': 5, 'K': 2, 'C': 1,
             'reasoning': 'Lay\'s: Maximum legacy. Bulk format = high B (decisions). Low C (autopilot snacking).'},
    'popchips': {'M': 3, 'E': 3, 'F': 3, 'B': 5, 'K': 2, 'C': 3,
                 'reasoning': 'Popchips: Moderate familiarity. "Better" positioning adds cognitive friction.'},
    'clif bar': {'M': 3, 'E': 3, 'F': 4, 'B': 1, 'K': 1, 'C': 2,
                 'reasoning': 'Clif Bar: Established energy bar brand (F=4). Unitized format = B=1, K=1. Some health consideration (C=2).'},
    'clif': {'M': 3, 'E': 3, 'F': 4, 'B': 1, 'K': 1, 'C': 2,
             'reasoning': 'Clif Bar: Established energy bar brand.'},
    'kind bar': {'M': 4, 'E': 3, 'F': 4, 'B': 1, 'K': 1, 'C': 2,
                 'reasoning': 'KIND: Well-established healthy snack (F=4). Unitized bar format.'},
    'kind': {'M': 4, 'E': 3, 'F': 4, 'B': 1, 'K': 1, 'C': 2,
             'reasoning': 'KIND: Well-established healthy snack.'},
    'rxbar': {'M': 3, 'E': 3, 'F': 3, 'B': 1, 'K': 1, 'C': 3,
              'reasoning': 'RXBAR: Growing brand, transparency messaging adds cognitive load.'},
    'joy days': {'M': 4, 'E': 4, 'F': 2, 'B': 1, 'K': 1, 'C': 3,
                 'reasoning': 'Joy Days: Newer frozen novelty brand (F=2). Unitized format. Some premium/health consideration.'},
    'joydays': {'M': 4, 'E': 4, 'F': 2, 'B': 1, 'K': 1, 'C': 3,
                'reasoning': 'Joy Days: Newer frozen novelty brand.'},
    'magnum': {'M': 5, 'E': 5, 'F': 4, 'B': 1, 'K': 1, 'C': 2,
               'reasoning': 'Magnum: Premium ice cream bar, established global brand (F=4). Unitized format.'},
    'mars': {'M': 4, 'E': 4, 'F': 5, 'B': 1, 'K': 1, 'C': 1,
             'reasoning': 'Mars: Iconic legacy brand (F=5). Autopilot purchase (C=1).'},
    'snickers': {'M': 5, 'E': 4, 'F': 5, 'B': 1, 'K': 1, 'C': 1,
                 'reasoning': 'Snickers: Maximum legacy (F=5). Autopilot purchase.'},
    'twix': {'M': 4, 'E': 4, 'F': 5, 'B': 1, 'K': 1, 'C': 1,
             'reasoning': 'Twix: Iconic candy bar. Maximum legacy.'},
    'talenti': {'M': 5, 'E': 5, 'F': 4, 'B': 4, 'K': 3, 'C': 2,
                'reasoning': 'Talenti: Premium gelato, well-known (F=4). Pint format = bulk denominator.'},
    'halo top': {'M': 3, 'E': 3, 'F': 3, 'B': 4, 'K': 3, 'C': 3,
                 'reasoning': 'Halo Top: Growing "better-for-you" ice cream. Cognitive for health positioning.'},
    'red bull': {'M': 3, 'E': 4, 'F': 5, 'B': 1, 'K': 1, 'C': 1,
                 'reasoning': 'Red Bull: Maximum legacy energy drink. Ritual can format.'},
    'monster': {'M': 3, 'E': 4, 'F': 4, 'B': 1, 'K': 1, 'C': 1,
                'reasoning': 'Monster: Established energy drink brand.'},
    'prime': {'M': 3, 'E': 3, 'F': 3, 'B': 1, 'K': 1, 'C': 4,
              'reasoning': 'Prime: Celebrity brand (Logan Paul/KSI). F=3 (new), C=4 (celebrity cognitive load).'},
    'feastables': {'M': 4, 'E': 4, 'F': 2, 'B': 2, 'K': 1, 'C': 5,
                   'reasoning': 'Feastables: MrBeast celebrity chocolate. Low F (new), High C (celebrity evaluation).'},
    'liquid death': {'M': 2, 'E': 3, 'F': 3, 'B': 1, 'K': 1, 'C': 3,
                     'reasoning': 'Liquid Death: Growing brand with edgy marketing. Some cognitive for branding novelty.'},
    'celsius': {'M': 3, 'E': 4, 'F': 3, 'B': 1, 'K': 1, 'C': 2,
                'reasoning': 'Celsius: Growing fitness energy drink brand.'},
}

# ═══════════════════════════════════════════════════════════════════════════════════════
# FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════════════

def levenshtein_distance(s1: str, s2: str) -> int:
    """Calculate edit distance between two strings."""
    if len(s1) < len(s2):
        return levenshtein_distance(s2, s1)
    if len(s2) == 0:
        return len(s1)
    
    previous_row = range(len(s2) + 1)
    for i, c1 in enumerate(s1):
        current_row = [i + 1]
        for j, c2 in enumerate(s2):
            insertions = previous_row[j + 1] + 1
            deletions = current_row[j] + 1
            substitutions = previous_row[j] + (c1 != c2)
            current_row.append(min(insertions, deletions, substitutions))
        previous_row = current_row
    
    return previous_row[-1]


def similarity_score(s1: str, s2: str) -> float:
    """Calculate similarity between two strings (0-1, higher is better)."""
    distance = levenshtein_distance(s1.lower(), s2.lower())
    max_len = max(len(s1), len(s2))
    if max_len == 0:
        return 1.0
    return 1 - (distance / max_len)


def find_similar_brands(query: str, threshold: float = 0.6) -> list:
    """Find brands similar to the query."""
    if not query:
        return []
    
    normalized = query.lower().strip()
    matches = []
    
    for brand_key in BRAND_DATABASE.keys():
        # Check exact/substring match first
        if brand_key in normalized or normalized in brand_key:
            return [(brand_key, 1.0)]  # Perfect match
        
        # Calculate similarity
        score = similarity_score(normalized, brand_key)
        
        # Also check without punctuation/spaces
        clean_query = ''.join(c for c in normalized if c.isalnum())
        clean_brand = ''.join(c for c in brand_key if c.isalnum())
        score2 = similarity_score(clean_query, clean_brand)
        
        best_score = max(score, score2)
        
        if best_score >= threshold:
            matches.append((brand_key, best_score))
    
    # Sort by similarity score (highest first)
    matches.sort(key=lambda x: x[1], reverse=True)
    return matches[:3]  # Return top 3 matches


def hunt_brand(brand_name: str) -> Tuple[Optional[str], Optional[str], bool, list]:
    """
    Hunt for brand in database with fuzzy matching.
    Returns: (archetype, description, is_ambiguous, similar_brands)
    """
    if not brand_name:
        return None, None, False, []
    
    normalized = brand_name.lower().strip()
    
    # Exact or substring match
    for brand_key, (archetype, desc) in BRAND_DATABASE.items():
        if brand_key in normalized or normalized in brand_key:
            return archetype, desc, False, []
    
    # Fuzzy matching
    similar = find_similar_brands(normalized)
    
    if similar:
        # If top match is very close (>0.8), auto-suggest correction
        if similar[0][1] >= 0.8:
            best_match = similar[0][0]
            archetype, desc = BRAND_DATABASE[best_match]
            return archetype, desc, False, similar
        else:
            # Return suggestions but no auto-match
            return None, None, False, similar
    
    return None, None, False, []

def calculate_s_score(m: int, e: int, f: int, b: int, k: int, c: int) -> float:
    return (m * e * f) / max(1, b * k * c)

def validate_rationale(text: str) -> bool:
    if not text:
        return False
    return sum(1 for c in text.strip() if c.isalnum()) >= 25

# ═══════════════════════════════════════════════════════════════════════════════════════
# MASSIVE FONT THEME
# ═══════════════════════════════════════════════════════════════════════════════════════

st.markdown("""
<style>
    @import url('https://fonts.googleapis.com/css2?family=Inter:wght@400;500;600;700;800;900&family=JetBrains+Mono:wght@400;500;600;700;800&display=swap');
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* DEEP VIOLET GLASSMORPHISM - WCAG AA ACCESSIBLE                                  */
    /* All text: Pure White (#FFFFFF) or Pale Lavender (#E9D5FF)                       */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    /* GLOBAL BACKGROUND - Deep Violet Gradient */
    .stApp {
        background: linear-gradient(135deg, #2e1065 0%, #4c1d95 100%) !important;
        background-attachment: fixed !important;
    }
    .main {
        background: transparent !important;
    }
    
    /* Watermark - very subtle */
    .main::before {
        content: 'ELBOW INTERFERENCE™';
        position: fixed; top: 50%; left: 50%;
        transform: translate(-50%, -50%) rotate(-25deg);
        font-size: 8rem; font-weight: 900;
        color: rgba(255, 255, 255, 0.02);
        pointer-events: none; z-index: 0;
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* GLOBAL TEXT RESET - FORCE HIGH CONTRAST                                         */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    /* EVERYTHING defaults to white */
    *, *::before, *::after {
        color: #FFFFFF;
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* SIDEBAR - Dark with WHITE text                                                  */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    section[data-testid="stSidebar"] { 
        background: rgba(15, 10, 35, 0.95) !important;
        border-right: 1px solid rgba(255, 255, 255, 0.15) !important;
        backdrop-filter: blur(12px) !important;
    }
    section[data-testid="stSidebar"] * {
        color: #FFFFFF !important;
    }
    section[data-testid="stSidebar"] h1,
    section[data-testid="stSidebar"] h2,
    section[data-testid="stSidebar"] h3,
    section[data-testid="stSidebar"] strong {
        color: #FFFFFF !important;
        text-shadow: 0 1px 2px rgba(0,0,0,0.3) !important;
    }
    section[data-testid="stSidebar"] p,
    section[data-testid="stSidebar"] span,
    section[data-testid="stSidebar"] label,
    section[data-testid="stSidebar"] .stCaption,
    section[data-testid="stSidebar"] [data-testid="stCaptionContainer"] {
        color: #E9D5FF !important;
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* TYPOGRAPHY - WCAG AA Compliant                                                  */
    /* Primary: #FFFFFF | Secondary: #E9D5FF | NEVER use grays                         */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    /* Headers - Bold White with text shadow */
    h1 { 
        font-size: 3.5rem !important; 
        font-weight: 900 !important; 
        color: #FFFFFF !important;
        text-shadow: 0 2px 4px rgba(0,0,0,0.3) !important;
    }
    h2 { 
        font-size: 2.5rem !important; 
        font-weight: 800 !important; 
        color: #FFFFFF !important; 
        margin-top: 2rem !important;
        text-shadow: 0 2px 4px rgba(0,0,0,0.3) !important;
    }
    h3 { 
        font-size: 1.8rem !important; 
        font-weight: 700 !important; 
        color: #FFFFFF !important;
        text-shadow: 0 1px 2px rgba(0,0,0,0.2) !important;
    }
    h4, h5 { 
        font-size: 1.4rem !important; 
        font-weight: 700 !important; 
        color: #E9D5FF !important;
    }
    
    /* Body text - Pale Lavender for readability */
    .main p, .main li, .main span { 
        font-size: 1.2rem !important; 
        line-height: 1.8 !important; 
        font-weight: 500 !important;
        color: #E9D5FF !important;
    }
    
    /* Labels - Pure White */
    .main label { 
        font-size: 1.2rem !important; 
        font-weight: 700 !important;
        color: #FFFFFF !important;
    }
    
    .main .stMarkdown p { color: #E9D5FF !important; }
    .main .stMarkdown strong { color: #FFFFFF !important; font-weight: 800 !important; }
    .main .stMarkdown li { color: #E9D5FF !important; }
    
    /* Captions - Pale Lavender (NOT muted gray) */
    small, .stCaption, [data-testid="stCaptionContainer"] { 
        font-size: 1rem !important; 
        font-weight: 500 !important; 
        color: #E9D5FF !important; 
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* GLASSMORPHISM CARDS - Darker backgrounds for contrast                           */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    /* Expander */
    [data-testid="stExpander"] {
        background: rgba(0, 0, 0, 0.25) !important;
        border: 1px solid rgba(255, 255, 255, 0.15) !important;
        backdrop-filter: blur(12px) !important;
        border-radius: 16px !important;
    }
    [data-testid="stExpander"] summary,
    [data-testid="stExpander"] summary *,
    .streamlit-expanderHeader,
    .streamlit-expanderHeader * {
        color: #FFFFFF !important;
        font-weight: 700 !important;
    }
    [data-testid="stExpander"] [data-testid="stMarkdownContainer"] *,
    .streamlit-expanderContent,
    .streamlit-expanderContent * { 
        color: #E9D5FF !important;
    }
    .streamlit-expanderContent strong { color: #FFFFFF !important; }
    
    /* Metrics - Glass Card */
    div[data-testid="stMetric"] { 
        background: rgba(0, 0, 0, 0.3) !important;
        border: 1px solid rgba(255, 255, 255, 0.15) !important;
        backdrop-filter: blur(12px) !important;
        border-radius: 16px !important; 
        padding: 24px !important; 
    }
    div[data-testid="stMetric"] label { 
        font-size: 1.1rem !important; 
        font-weight: 700 !important; 
        text-transform: uppercase !important; 
        letter-spacing: 0.1em !important; 
        color: #E9D5FF !important;
    }
    div[data-testid="stMetric"] div[data-testid="stMetricValue"] { 
        font-family: 'JetBrains Mono', monospace !important; 
        font-size: 3rem !important; 
        font-weight: 800 !important; 
        color: #10b981 !important;
        text-shadow: 0 0 20px rgba(16, 185, 129, 0.5) !important;
    }
    div[data-testid="stMetric"] div[data-testid="stMetricDelta"] { 
        font-size: 1.2rem !important; 
        font-weight: 700 !important;
        color: #ec4899 !important;
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* INPUTS - DARK backgrounds (bg-black/30), WHITE text, white/40 placeholders      */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    /* Text Inputs */
    .stTextInput label { 
        font-size: 1.3rem !important; 
        font-weight: 700 !important; 
        color: #FFFFFF !important;
        text-shadow: 0 1px 2px rgba(0,0,0,0.2) !important;
    }
    .stTextInput input { 
        font-size: 1.2rem !important; 
        padding: 16px 20px !important; 
        background: rgba(0, 0, 0, 0.4) !important;
        border: 1px solid rgba(255, 255, 255, 0.2) !important;
        border-radius: 12px !important;
        color: #FFFFFF !important;
    }
    .stTextInput input::placeholder { 
        color: rgba(255, 255, 255, 0.4) !important;
        opacity: 1 !important;
    }
    .stTextInput input:focus { 
        border-color: #10b981 !important;
        box-shadow: 0 0 0 3px rgba(16, 185, 129, 0.3) !important;
        outline: none !important;
    }
    
    /* Number Inputs */
    .stNumberInput label { 
        font-size: 1.1rem !important; 
        font-weight: 700 !important; 
        color: #FFFFFF !important;
    }
    .stNumberInput > div { 
        background: rgba(0, 0, 0, 0.3) !important;
        border-radius: 12px !important;
        padding: 4px !important;
        border: 1px solid rgba(255, 255, 255, 0.15) !important;
    }
    .stNumberInput input { 
        font-size: 1.5rem !important; 
        font-weight: 800 !important;
        padding: 8px 4px !important; 
        text-align: center !important;
        width: 60px !important;
        border: 1px solid rgba(255, 255, 255, 0.25) !important;
        border-radius: 8px !important;
        background: rgba(0, 0, 0, 0.4) !important;
        color: #FFFFFF !important;
    }
    .stNumberInput button {
        font-size: 1.1rem !important;
        min-width: 32px !important;
        height: 32px !important;
        background: #10b981 !important;
        color: #FFFFFF !important;
        border-radius: 6px !important;
        border: none !important;
        margin: 0 4px !important;
    }
    .stNumberInput button:hover { background: #059669 !important; }
    .stNumberInput button svg { stroke: #FFFFFF !important; fill: #FFFFFF !important; }
    .stNumberInput button * { color: #FFFFFF !important; -webkit-text-fill-color: #FFFFFF !important; }
    
    /* Text Areas */
    .stTextArea label { 
        font-size: 1.2rem !important; 
        font-weight: 700 !important; 
        color: #FFFFFF !important;
    }
    .stTextArea textarea { 
        font-size: 1.1rem !important; 
        padding: 16px !important; 
        background: rgba(0, 0, 0, 0.4) !important;
        border: 1px solid rgba(255, 255, 255, 0.2) !important;
        border-radius: 12px !important; 
        color: #FFFFFF !important;
        min-height: 120px !important; 
    }
    .stTextArea textarea::placeholder {
        color: rgba(255, 255, 255, 0.4) !important;
    }
    
    /* Selectbox */
    .stSelectbox label { 
        font-size: 1.2rem !important; 
        font-weight: 700 !important; 
        color: #FFFFFF !important; 
    }
    .stSelectbox > div > div { 
        background: rgba(0, 0, 0, 0.4) !important;
        color: #FFFFFF !important;
        border: 1px solid rgba(255, 255, 255, 0.2) !important;
        border-radius: 12px !important;
        min-height: 48px !important;
    }
    .stSelectbox [data-baseweb="select"] { background: transparent !important; }
    .stSelectbox [data-baseweb="select"] * { color: #FFFFFF !important; }
    .stSelectbox span, .stSelectbox div { color: #FFFFFF !important; }
    
    /* Dropdown menu - DARK background */
    div[data-baseweb="popover"], div[data-baseweb="popover"] * { 
        background: rgba(15, 10, 35, 0.98) !important; 
        color: #FFFFFF !important;
        backdrop-filter: blur(12px) !important;
    }
    div[data-baseweb="popover"] li { color: #FFFFFF !important; }
    div[data-baseweb="popover"] li:hover { background: rgba(16, 185, 129, 0.4) !important; }
    div[data-baseweb="menu"], div[data-baseweb="menu"] * { 
        background: rgba(15, 10, 35, 0.98) !important; 
        color: #FFFFFF !important;
    }
    [data-baseweb="list"], [data-baseweb="list"] * { 
        background: rgba(15, 10, 35, 0.98) !important; 
        color: #FFFFFF !important;
    }
    [data-baseweb="list-item"], [data-baseweb="list-item"] * { 
        background: transparent !important; 
        color: #FFFFFF !important;
    }
    [data-baseweb="list-item"]:hover { background: rgba(16, 185, 129, 0.4) !important; }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* ALERTS - High contrast text                                                     */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    .stAlert { 
        background: rgba(0, 0, 0, 0.3) !important;
        border: 1px solid rgba(255, 255, 255, 0.15) !important;
        backdrop-filter: blur(12px) !important;
        border-radius: 12px !important;
        padding: 16px 20px !important;
    }
    .stAlert * { color: #FFFFFF !important; }
    .stAlert p { color: #FFFFFF !important; font-weight: 600 !important; }
    
    /* Success - Mint */
    [data-testid="stAlert"]:has(svg[data-testid="stNotificationContentSuccess"]) {
        background: rgba(16, 185, 129, 0.2) !important;
        border-left: 4px solid #10b981 !important;
    }
    
    /* Error */
    [data-testid="stAlert"]:has(svg[data-testid="stNotificationContentError"]) {
        background: rgba(239, 68, 68, 0.2) !important;
        border-left: 4px solid #ef4444 !important;
    }
    
    /* Info */
    [data-testid="stAlert"]:has(svg[data-testid="stNotificationContentInfo"]) {
        background: rgba(59, 130, 246, 0.2) !important;
        border-left: 4px solid #3b82f6 !important;
    }
    
    /* Warning */
    [data-testid="stAlert"]:has(svg[data-testid="stNotificationContentWarning"]) {
        background: rgba(236, 72, 153, 0.2) !important;
        border-left: 4px solid #ec4899 !important;
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* BUTTONS - Bright accents with white text                                        */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    .stButton button { 
        font-size: 1.1rem !important; 
        font-weight: 700 !important; 
        padding: 14px 24px !important; 
        border-radius: 12px !important;
        background: linear-gradient(135deg, #10b981 0%, #059669 100%) !important;
        color: #FFFFFF !important;
        border: none !important;
        transition: all 0.2s ease !important;
        text-shadow: 0 1px 2px rgba(0,0,0,0.2) !important;
    }
    .stButton button:hover {
        transform: translateY(-2px) !important;
        box-shadow: 0 8px 20px rgba(16, 185, 129, 0.4) !important;
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* OTHER COMPONENTS                                                                */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    /* Dividers */
    hr { 
        margin: 24px 0 !important; 
        border: none !important;
        border-top: 1px solid rgba(255, 255, 255, 0.1) !important;
    }
    
    /* DataFrames */
    .stDataFrame { 
        background: rgba(0, 0, 0, 0.2) !important;
        border-radius: 12px !important;
    }
    .stDataFrame * { color: #FFFFFF !important; }
    
    /* Code blocks */
    code { 
        font-size: 1.1rem !important; 
        font-weight: 600 !important; 
        background: rgba(0, 0, 0, 0.3) !important; 
        color: #10b981 !important;
        padding: 4px 10px !important; 
        border-radius: 6px !important;
        border: 1px solid rgba(255, 255, 255, 0.1) !important;
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* FILE UPLOADER - Glass Card Style (NO WHITE BACKGROUNDS)                         */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    [data-testid="stFileUploader"],
    [data-testid="stFileUploader"] > div,
    [data-testid="stFileUploader"] > div > div,
    [data-testid="stFileUploader"] section,
    [data-testid="stFileUploadDropzone"],
    [data-testid="stFileUploadDropzone"] > div {
        background: rgba(0, 0, 0, 0.4) !important;
        background-color: rgba(0, 0, 0, 0.4) !important;
        border-radius: 12px !important;
    }
    
    /* The actual drop zone */
    [data-testid="stFileUploadDropzone"] {
        background: rgba(15, 10, 35, 0.6) !important;
        border: 2px dashed rgba(255, 255, 255, 0.25) !important;
        border-radius: 12px !important;
        padding: 24px !important;
    }
    [data-testid="stFileUploadDropzone"]:hover {
        border-color: #10b981 !important;
        background: rgba(16, 185, 129, 0.1) !important;
    }
    
    /* All text in file uploader - WHITE */
    [data-testid="stFileUploader"] *,
    [data-testid="stFileUploadDropzone"] * { 
        color: #FFFFFF !important; 
        background-color: transparent !important;
    }
    [data-testid="stFileUploader"] span,
    [data-testid="stFileUploader"] p,
    [data-testid="stFileUploader"] small { 
        color: #E9D5FF !important; 
    }
    
    /* Browse button */
    [data-testid="stFileUploader"] button,
    [data-testid="stFileUploadDropzone"] button { 
        background: linear-gradient(135deg, #10b981 0%, #059669 100%) !important;
        color: #FFFFFF !important;
        border: none !important;
        padding: 8px 16px !important;
        border-radius: 8px !important;
    }
    [data-testid="stFileUploader"] button:hover {
        background: #10b981 !important;
        box-shadow: 0 4px 12px rgba(16, 185, 129, 0.3) !important;
    }
    
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    /* CHARTS - Transparent Background, White Labels                                   */
    /* ═══════════════════════════════════════════════════════════════════════════════ */
    
    /* Remove white backgrounds from chart containers */
    [data-testid="stVegaLiteChart"],
    [data-testid="stVegaLiteChart"] > div,
    [data-testid="stArrowVegaLiteChart"],
    [data-testid="stArrowVegaLiteChart"] > div,
    .vega-embed,
    .vega-embed > div,
    .vega-embed canvas,
    .marks {
        background: transparent !important;
        background-color: transparent !important;
    }
    
    /* Chart wrapper - glass card */
    [data-testid="stVegaLiteChart"],
    [data-testid="stArrowVegaLiteChart"] {
        background: rgba(0, 0, 0, 0.3) !important;
        border: 1px solid rgba(255, 255, 255, 0.1) !important;
        border-radius: 16px !important;
        padding: 16px !important;
    }
    
    /* Force Vega-Lite text to white */
    .vega-embed text,
    .vega-embed .role-axis-label,
    .vega-embed .role-legend-label,
    .vega-embed .role-title,
    .vega-embed .mark-text text {
        fill: #FFFFFF !important;
        color: #FFFFFF !important;
    }
    
    /* Axis lines and grid - light */
    .vega-embed .role-axis line,
    .vega-embed .role-axis path,
    .vega-embed line.role-grid {
        stroke: rgba(255, 255, 255, 0.2) !important;
    }
    
    /* Legend text */
    .vega-embed .role-legend text {
        fill: #E9D5FF !important;
    }
    
    /* Icons - White or Mint */
    svg { color: #FFFFFF !important; }
    svg:not(.vega-embed svg) { stroke: #FFFFFF !important; }
    .stAlert svg { color: #10b981 !important; stroke: #10b981 !important; }
    
    /* LOCKED STATE - Alert Red */
    .score-locked {
        background: rgba(239, 68, 68, 0.2) !important;
        border: 2px solid #ef4444 !important;
        animation: pulse 1.5s infinite;
    }
    @keyframes pulse {
        0%, 100% { opacity: 1; box-shadow: 0 0 0 0 rgba(239, 68, 68, 0.4); }
        50% { opacity: 0.8; box-shadow: 0 0 20px 5px rgba(239, 68, 68, 0.2); }
    }
    
    /* Hide Streamlit branding */
    #MainMenu, footer, header { visibility: hidden; }
</style>
<script>document.addEventListener('contextmenu', e => e.preventDefault());</script>
""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════════════════════════════
# SESSION STATE
# ═══════════════════════════════════════════════════════════════════════════════════════

if "brand1_archetype" not in st.session_state: st.session_state.brand1_archetype = None
if "brand2_archetype" not in st.session_state: st.session_state.brand2_archetype = None
if "rationale_b1" not in st.session_state: st.session_state.rationale_b1 = ""
if "rationale_b2" not in st.session_state: st.session_state.rationale_b2 = ""
if "ai_scores_b1" not in st.session_state: st.session_state.ai_scores_b1 = None
if "ai_scores_b2" not in st.session_state: st.session_state.ai_scores_b2 = None
if "last_brand1" not in st.session_state: st.session_state.last_brand1 = ""
if "last_brand2" not in st.session_state: st.session_state.last_brand2 = ""

# Priors context state (behavioral adjustments)
if "use_priors" not in st.session_state: st.session_state.use_priors = True
if "cohort" not in st.session_state: st.session_state.cohort = "mixed"
if "occasion" not in st.session_state: st.session_state.occasion = "evening"
if "macro_stress" not in st.session_state: st.session_state.macro_stress = True
if "promo_frequency_b1" not in st.session_state: st.session_state.promo_frequency_b1 = 0.0
if "promo_depth_b1" not in st.session_state: st.session_state.promo_depth_b1 = 0.0
if "promo_frequency_b2" not in st.session_state: st.session_state.promo_frequency_b2 = 0.0
if "promo_depth_b2" not in st.session_state: st.session_state.promo_depth_b2 = 0.0
if "ups_pw_13_b1" not in st.session_state: st.session_state.ups_pw_13_b1 = None
if "ups_pw_26_b1" not in st.session_state: st.session_state.ups_pw_26_b1 = None
if "ups_pw_13_b2" not in st.session_state: st.session_state.ups_pw_13_b2 = None
if "ups_pw_26_b2" not in st.session_state: st.session_state.ups_pw_26_b2 = None
if "adjustments_log_b1" not in st.session_state: st.session_state.adjustments_log_b1 = []
if "adjustments_log_b2" not in st.session_state: st.session_state.adjustments_log_b2 = []

# Check if Gemini is configured
_test_model = get_gemini_model()
AI_ENABLED = _test_model is not None
print(f"[STARTUP] AI_ENABLED = {AI_ENABLED}")

# ═══════════════════════════════════════════════════════════════════════════════════════
# SIDEBAR
# ═══════════════════════════════════════════════════════════════════════════════════════

with st.sidebar:
    st.header("◈ Elbow Interference™")
    
    # Status indicator
    if AI_ENABLED:
        st.success("System Active")
    else:
        st.warning("Manual Mode")
    
    st.divider()
    
    # RESTART
    if st.button("🔄 Restart Analysis", use_container_width=True):
        st.session_state.ai_scores_b1 = None
        st.session_state.ai_scores_b2 = None
        st.session_state.last_brand1 = ""
        st.session_state.last_brand2 = ""
        st.session_state['final_b1'] = {}
        st.session_state['final_b2'] = {}
        st.session_state.last_memo = None
        st.rerun()
    
    st.divider()
    
    # Data Upload - Multiple file types
    st.subheader("📊 Data Upload")
    uploaded_file = st.file_uploader(
        "Upload Data", 
        type=["csv", "xlsx", "xls", "pdf", "docx", "doc", "pptx", "ppt", "txt"],
        help=tip("upload")
    )
    
    if uploaded_file:
        file_ext = uploaded_file.name.split('.')[-1].lower()
        file_size_mb = uploaded_file.size / (1024 * 1024)
        
        # Show processing status
        st.info(f"📄 **{uploaded_file.name}** ({file_size_mb:.1f} MB)")
        
        try:
            # CSV files
            if file_ext == 'csv':
                with st.spinner("⏳ Processing CSV..."):
                    df = pd.read_csv(uploaded_file)
                    st.session_state['uploaded_data'] = df
                    st.session_state['uploaded_text'] = None
                    st.session_state['file_loaded'] = True
                st.success(f"✅ **CSV loaded:** {len(df)} rows × {len(df.columns)} columns")
                
                # Show column preview
                with st.expander("📋 Data Preview"):
                    try:
                        st.dataframe(df.head(10), use_container_width=True)
                    except:
                        st.table(df.head(10))
                
                # Time series detection for CSV
                time_cols = [col for col in df.columns if any(x in col.lower() for x in ['52', '26', '13', '4', 'week', 'wk', 'latest'])]
                if time_cols:
                    selected_period = st.selectbox("Time Period", ["All Provided"] + time_cols, key="data_period")
                    st.session_state['selected_period'] = selected_period
            
            # Excel files
            elif file_ext in ['xlsx', 'xls']:
                with st.spinner("⏳ Processing Excel file... (large files may take a moment)"):
                    # Try to read with openpyxl for xlsx
                    if file_ext == 'xlsx':
                        df = pd.read_excel(uploaded_file, engine='openpyxl')
                    else:
                        df = pd.read_excel(uploaded_file)
                    
                    # Clean up mixed data types - convert all columns to strings for display
                    df_clean = df.copy()
                    for col in df_clean.columns:
                        df_clean[col] = df_clean[col].astype(str)
                    
                    st.session_state['uploaded_data'] = df
                    st.session_state['uploaded_text'] = None
                    st.session_state['file_loaded'] = True
                    
                st.success(f"✅ **Excel loaded:** {len(df)} rows × {len(df.columns)} columns")
                st.caption(f"Columns: {', '.join(df.columns[:5].tolist())}{'...' if len(df.columns) > 5 else ''}")
                
                # Show column preview with cleaned data
                with st.expander("📋 Data Preview (first 10 rows)"):
                    try:
                        st.dataframe(df_clean.head(10), use_container_width=True)
                    except:
                        st.table(df.head(10))
                
                # Time series detection for Excel
                time_cols = [col for col in df.columns if any(x in str(col).lower() for x in ['52', '26', '13', '4', 'week', 'wk', 'latest'])]
                if time_cols:
                    selected_period = st.selectbox("Time Period", ["All Provided"] + time_cols, key="data_period")
                    st.session_state['selected_period'] = selected_period
            
            # PDF files
            elif file_ext == 'pdf':
                try:
                    import PyPDF2
                    with st.spinner("⏳ Extracting text from PDF..."):
                        pdf_reader = PyPDF2.PdfReader(uploaded_file)
                        text_content = ""
                        for page in pdf_reader.pages:
                            text_content += page.extract_text() + "\n"
                        st.session_state['uploaded_text'] = text_content
                        st.session_state['uploaded_data'] = None
                    st.success(f"✅ PDF extracted: {len(pdf_reader.pages)} pages, {len(text_content)} characters")
                    with st.expander("📋 Preview extracted text"):
                        st.text(text_content[:3000] + "..." if len(text_content) > 3000 else text_content)
                except ImportError:
                    st.error("📦 Missing library: `pip install PyPDF2`")
                except Exception as pdf_err:
                    st.error(f"PDF Error: {str(pdf_err)[:100]}")
            
            # Word documents
            elif file_ext in ['docx', 'doc']:
                try:
                    from docx import Document
                    with st.spinner("⏳ Extracting text from Word document..."):
                        doc = Document(uploaded_file)
                        text_content = "\n".join([para.text for para in doc.paragraphs])
                        st.session_state['uploaded_text'] = text_content
                        st.session_state['uploaded_data'] = None
                    st.success(f"✅ Word extracted: {len(doc.paragraphs)} paragraphs")
                    with st.expander("📋 Preview extracted text"):
                        st.text(text_content[:3000] + "..." if len(text_content) > 3000 else text_content)
                except ImportError:
                    st.error("📦 Missing library: `pip install python-docx`")
                except Exception as doc_err:
                    st.error(f"Word Error: {str(doc_err)[:100]}")
            
            # PowerPoint files
            elif file_ext in ['pptx', 'ppt']:
                try:
                    from pptx import Presentation
                    with st.spinner("⏳ Extracting text from PowerPoint..."):
                        prs = Presentation(uploaded_file)
                        text_content = ""
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text_content += shape.text + "\n"
                        st.session_state['uploaded_text'] = text_content
                        st.session_state['uploaded_data'] = None
                    st.success(f"✅ PowerPoint extracted: {len(prs.slides)} slides")
                    with st.expander("📋 Preview extracted text"):
                        st.text(text_content[:3000] + "..." if len(text_content) > 3000 else text_content)
                except ImportError:
                    st.error("📦 Missing library: `pip install python-pptx`")
                except Exception as ppt_err:
                    st.error(f"PowerPoint Error: {str(ppt_err)[:100]}")
            
            # Plain text files
            elif file_ext == 'txt':
                with st.spinner("⏳ Loading text file..."):
                    text_content = uploaded_file.read().decode('utf-8')
                    st.session_state['uploaded_text'] = text_content
                    st.session_state['uploaded_data'] = None
                st.success(f"✅ Text loaded: {len(text_content)} characters")
                with st.expander("📋 Preview text"):
                    st.text(text_content[:3000] + "..." if len(text_content) > 3000 else text_content)
            
            else:
                st.warning(f"⚠️ Unsupported file type: {file_ext}")
                
        except Exception as e:
            st.error(f"❌ Error reading file: {str(e)}")
            st.session_state['file_loaded'] = False
            import traceback
            with st.expander("🔧 Error Details"):
                st.code(traceback.format_exc())
    
    st.divider()
    
    # ═══════════════════════════════════════════════════════════════════════════════════════
    # ADVANCED SETTINGS - BEHAVIORAL PRIORS
    # ═══════════════════════════════════════════════════════════════════════════════════════
    
    with st.expander("⚙️ Advanced Settings", expanded=False):
        st.markdown("**Behavioral Priors**")
        
        # Feature flag toggle
        st.session_state.use_priors = st.checkbox(
            "Apply Behavioral Priors",
            value=st.session_state.use_priors,
            help=tip("priors_toggle")
        )
        
        if st.session_state.use_priors and PRIORS_AVAILABLE:
            st.caption("Context affects scoring adjustments:")
            
            st.session_state.cohort = st.selectbox(
                "Target Cohort",
                options=["younger", "mixed", "older"],
                index=["younger", "mixed", "older"].index(st.session_state.cohort),
                help=tip("cohort")
            )
            
            st.session_state.occasion = st.selectbox(
                "Primary Occasion",
                options=["evening", "late_night", "daytime", "on_the_go"],
                index=["evening", "late_night", "daytime", "on_the_go"].index(st.session_state.occasion),
                help=tip("priors_occasion")
            )
            
            st.session_state.macro_stress = st.checkbox(
                "Macro Stress Active",
                value=st.session_state.macro_stress,
                help=tip("macro_stress")
            )
            
            st.divider()
            st.caption("**Promo Reliance (Brand 1)**")
            st.session_state.promo_frequency_b1 = st.slider(
                "Promo Frequency B1",
                0.0, 1.0, st.session_state.promo_frequency_b1,
                help=tip("promo_frequency")
            )
            st.session_state.promo_depth_b1 = st.slider(
                "Promo Depth B1",
                0.0, 1.0, st.session_state.promo_depth_b1,
                help=tip("promo_depth")
            )
            
            st.caption("**Promo Reliance (Brand 2)**")
            st.session_state.promo_frequency_b2 = st.slider(
                "Promo Frequency B2",
                0.0, 1.0, st.session_state.promo_frequency_b2,
                help=tip("promo_frequency")
            )
            st.session_state.promo_depth_b2 = st.slider(
                "Promo Depth B2",
                0.0, 1.0, st.session_state.promo_depth_b2,
                help=tip("promo_depth")
            )
        elif not PRIORS_AVAILABLE:
            st.warning("Priors module not loaded")
    
    st.divider()
    
    # FORMAT REFERENCE
    st.subheader("Format Reference")
    
    st.success("SINGLE-SERVE (Low Friction): Bars, bites, pouches. B=1, K=1")
    st.error("MULTI-SERVE (High Friction): Pints, tubs, bags. B=4, K=3")
    st.info("RITUAL DRINKS (Low Friction): Cans, bottles. F=5, B=1, K=1")
    
    st.divider()
    
    st.subheader("The Equation")
    st.markdown("""
<div style="background: rgba(0, 0, 0, 0.3); padding: 12px 16px; border-radius: 8px; border: 1px solid rgba(255, 255, 255, 0.1);">
    <code style="color: #10b981 !important; font-size: 1.1rem; font-weight: 600;">S = (M×E×F) ÷ (B×K×C)</code>
</div>
    """, unsafe_allow_html=True)
    st.caption("Value Delivered ÷ Cost Extracted = Satisfaction")
    
    st.divider()
    
    st.caption("© Russell Barnett 2026")

# ═══════════════════════════════════════════════════════════════════════════════════════
# MAIN HEADER
# ═══════════════════════════════════════════════════════════════════════════════════════

st.markdown("# Elbow Interference Evaluator™")

# Compact Equation Banner - Glassmorphism
st.markdown("""
<div style="background: rgba(91, 33, 182, 0.35); backdrop-filter: blur(12px); border: 1px solid rgba(255, 255, 255, 0.15); padding: 24px 40px; border-radius: 16px; display: flex; align-items: center; justify-content: space-between; margin: 16px 0 24px 0;">
    <div style="font-family: 'JetBrains Mono', monospace; font-size: 1.8rem; font-weight: 700; color: #10b981 !important;">S = (M × E × F) ÷ (B × K × C)</div>
    <div style="font-size: 1rem; color: #c4b5fd !important; text-align: right;">S = Satisfaction (enables persistence)</div>
</div>
""", unsafe_allow_html=True)

# ═══════════════════════════════════════════════════════════════════════════════════════
# BRAND DISCOVERY - CLEAN: BRAND + FORMAT SEPARATE
# ═══════════════════════════════════════════════════════════════════════════════════════

st.markdown("## 🔍 Brand Discovery")

# Educational context
with st.expander("📖 How Scoring Works (Elbow Interference Theory™)", expanded=False):
    st.markdown("""
    **S = Satisfaction.** Value Delivered ÷ Cost Extracted = Satisfaction. Satisfaction enables persistence over time, but persistence is not the equation itself.

    **The Key Insight:** Products don't win by maximizing pleasure—they win by minimizing friction.
    
    **F (Familiarity)** — Does the consumer buy on autopilot?
    - **5** = Iconic legacy (Coke, Oreo, Lay's) — decades of habit
    - **3** = Growing brand or celebrity launch — still building ritual
    - **1** = Unknown/new — requires discovery
    
    **C (Cognitive)** — Does the brand make consumers THINK?
    - **1-2** = Autopilot purchase — grab without thinking
    - **3** = Some consideration — premium or health-adjacent
    - **4-5** = High cognitive load — celebrity brands, complex value props
    
    **Celebrity brands = Low F + High C** (the celebrity is familiar, but the PRODUCT is new and requires evaluation)
    
    **Format determines B and K:**
    - **Single-serve** = B=1, K=1 (package ends the occasion)
    - **Multi-serve** = B=4, K=3 (you decide when to stop)
    """)

# Default category when not set in session (used for priors context)
DEFAULT_CATEGORY = "ice_cream"
CATEGORY_OPTIONS = ["ice_cream", "chips", "candy", "soda", "energy", "yogurt", "bars", "other"]

# Format options - covers ALL orally consumed CPG
FORMAT_LIST = [
    "Pint / Tub (multi-serve)",
    "Bag / Pouch (multi-serve)", 
    "Box / Carton (multi-serve)",
    "Single-Serve Bar / Novelty",
    "Can / Bottle (single)",
]
FORMAT_TO_ARCHETYPE = {
    "Pint / Tub (multi-serve)": "bulk",
    "Bag / Pouch (multi-serve)": "bulk",
    "Box / Carton (multi-serve)": "bulk",
    "Single-Serve Bar / Novelty": "unitized",
    "Can / Bottle (single)": "ritual",
}

# Optional occasion context (free-text, separate from priors occasion enum)
occasion_free_text = st.text_input(
    "Occasion (optional)",
    placeholder="e.g., snacking, dessert, breakfast...",
    key="occasion_free_text",
    help=tip("occasion_free_text"),
)

st.markdown("---")

h1, h2 = st.columns(2)

with h1:
    st.markdown("### 🟢 BRAND 1 (Incumbent)")
    brand1_name = st.text_input(
        "Brand Name",
        placeholder="e.g., Serendipity, Dr. Bombay...",
        key="b1",
        help=tip("brand_name"),
    )
    brand1_format = st.selectbox(
        "Product Format", 
        options=["Pint / Tub (multi-serve)", "Bag / Pouch (multi-serve)", "Box / Carton (multi-serve)", "Single-Serve Bar / Novelty", "Can / Bottle (single)"],
        index=0, 
        key="b1_format",
        help=tip("format"),
    )
    st.caption(f"**Selected: {brand1_format}**")
    st.session_state.brand1_archetype = FORMAT_TO_ARCHETYPE[brand1_format]
    st.selectbox(
        "Category",
        options=CATEGORY_OPTIONS,
        index=0,
        key="brand1_category",
    )
    
    # Get format-based defaults
    arch1 = ARCHETYPES[FORMAT_TO_ARCHETYPE[brand1_format]]
    
    if brand1_name:
        # Normalize brand name for lookup (uses global KNOWN_BRANDS)
        brand_key = brand1_name.lower().strip()
        known_scores = KNOWN_BRANDS.get(brand_key)
        
        # Try known database first, then AI
        if known_scores:
            st.session_state.ai_scores_b1 = known_scores
            st.success(f"✓ **{brand1_name}** — F={known_scores['F']}, C={known_scores['C']}")
            with st.expander("Analysis", expanded=False):
                st.write(known_scores.get('reasoning', ''))
        elif AI_ENABLED and (st.session_state.ai_scores_b1 is None or st.session_state.last_brand1 != brand1_name):
            st.session_state.last_brand1 = brand1_name
            with st.spinner(f"Analyzing {brand1_name}..."):
                try:
                    ai_result = analyze_brand_with_ai(brand1_name)
                    if ai_result and "error" not in ai_result:
                        st.session_state.ai_scores_b1 = ai_result
                        st.success(f"✓ **{brand1_name}** — F={ai_result.get('F')}, C={ai_result.get('C')}")
                    else:
                        st.warning(f"⚠️ **{brand1_name}** not recognized — set scores manually")
                except Exception as e:
                    st.warning(f"⚠️ Set scores manually for {brand1_name}")
                    print(f"[AI ERROR B1]: {e}")
        elif st.session_state.ai_scores_b1:
            ai1 = st.session_state.ai_scores_b1
            st.success(f"✓ **{brand1_name}** — F={ai1.get('F')}, C={ai1.get('C')}")
        else:
            st.warning(f"⚠️ **{brand1_name}** — set scores manually below")
        
        ai1 = st.session_state.ai_scores_b1
        
        # Store original AI scores for comparison
        orig_M1 = ai1.get('M', 4) if ai1 else 4
        orig_E1 = ai1.get('E', 4) if ai1 else 4
        orig_F1 = ai1.get('F', 3) if ai1 else 3
        orig_B1 = ai1.get('B', arch1.B) if ai1 else arch1.B
        orig_K1 = ai1.get('K', arch1.K) if ai1 else arch1.K
        orig_C1 = ai1.get('C', 3) if ai1 else 3
        
        # === ALL 6 SCORE CONTROLS ===
        st.markdown("**ALL 6 METRICS** (adjust as needed)")
        
        # Row 1: M, E, F
        row1_c1, row1_c2, row1_c3 = st.columns(3)
        with row1_c1:
            b1_manual_M = st.number_input("M", min_value=1, max_value=5, value=orig_M1, key="b1_M", help=tip("M"))
            st.caption("Mouthfeel")
        with row1_c2:
            b1_manual_E = st.number_input("E", min_value=1, max_value=5, value=orig_E1, key="b1_E", help=tip("E"))
            st.caption("Emotion")
        with row1_c3:
            b1_manual_F = st.number_input("F ⭐", min_value=1, max_value=5, value=orig_F1, key="b1_F", help=tip("F"))
            st.caption("Familiarity")
        
        # Row 2: B, K, C
        row2_c1, row2_c2, row2_c3 = st.columns(3)
        with row2_c1:
            b1_manual_B = st.number_input("B", min_value=1, max_value=5, value=orig_B1, key="b1_B", help=tip("B"))
            st.caption("Bites")
        with row2_c2:
            b1_manual_K = st.number_input("K", min_value=1, max_value=5, value=orig_K1, key="b1_K", help=tip("K"))
            st.caption("Kinetic")
        with row2_c3:
            b1_manual_C = st.number_input("C ⭐", min_value=1, max_value=5, value=orig_C1, key="b1_C", help=tip("C"))
            st.caption("Cognitive")
        
        # Check for overrides - each needs individual justification
        b1_justifications = {}
        b1_missing_justifications = []
        
        if b1_manual_M != orig_M1:
            st.warning(f"⚠️ **M Override**: {orig_M1}→{b1_manual_M}")
            b1_justifications['M'] = st.text_input(f"Why M={b1_manual_M}?", key="b1_just_M", placeholder="Explain Mouthfeel change...", help=tip("override_rationale"))
            if len(b1_justifications['M'].strip()) < 10: b1_missing_justifications.append('M')
        
        if b1_manual_E != orig_E1:
            st.warning(f"⚠️ **E Override**: {orig_E1}→{b1_manual_E}")
            b1_justifications['E'] = st.text_input(f"Why E={b1_manual_E}?", key="b1_just_E", placeholder="Explain Emotion change...", help=tip("override_rationale"))
            if len(b1_justifications['E'].strip()) < 10: b1_missing_justifications.append('E')
        
        if b1_manual_F != orig_F1:
            st.warning(f"⚠️ **F Override**: {orig_F1}→{b1_manual_F}")
            b1_justifications['F'] = st.text_input(f"Why F={b1_manual_F}?", key="b1_just_F", placeholder="Explain Familiarity change...", help=tip("override_rationale"))
            if len(b1_justifications['F'].strip()) < 10: b1_missing_justifications.append('F')
        
        if b1_manual_B != orig_B1:
            st.warning(f"⚠️ **B Override**: {orig_B1}→{b1_manual_B}")
            b1_justifications['B'] = st.text_input(f"Why B={b1_manual_B}?", key="b1_just_B", placeholder="Explain Bites change...", help=tip("override_rationale"))
            if len(b1_justifications['B'].strip()) < 10: b1_missing_justifications.append('B')
        
        if b1_manual_K != orig_K1:
            st.warning(f"⚠️ **K Override**: {orig_K1}→{b1_manual_K}")
            b1_justifications['K'] = st.text_input(f"Why K={b1_manual_K}?", key="b1_just_K", placeholder="Explain Kinetic change...", help=tip("override_rationale"))
            if len(b1_justifications['K'].strip()) < 10: b1_missing_justifications.append('K')
        
        if b1_manual_C != orig_C1:
            st.warning(f"⚠️ **C Override**: {orig_C1}→{b1_manual_C}")
            b1_justifications['C'] = st.text_input(f"Why C={b1_manual_C}?", key="b1_just_C", placeholder="Explain Cognitive change...", help=tip("override_rationale"))
            if len(b1_justifications['C'].strip()) < 10: b1_missing_justifications.append('C')
        
        if b1_missing_justifications:
            st.error(f"❌ Justify: {', '.join(b1_missing_justifications)} (min 10 chars each)")
            st.session_state['b1_locked'] = True
        elif b1_justifications:
            st.success("✓ All overrides justified")
            st.session_state['b1_locked'] = False
        else:
            st.session_state['b1_locked'] = False
        
        b1_justification = " | ".join([f"{k}: {v}" for k, v in b1_justifications.items() if v])
        
        # Store final values
        final_scores_b1 = {
            'M': b1_manual_M, 'E': b1_manual_E, 'F': b1_manual_F,
            'B': b1_manual_B, 'K': b1_manual_K, 'C': b1_manual_C,
            'archetype': FORMAT_TO_ARCHETYPE[brand1_format],
            'reasoning': ai1.get('reasoning', 'Manual scoring') if ai1 else 'Manual scoring',
            'override_justification': b1_justification if b1_justifications else ''
        }
        st.session_state['final_b1'] = final_scores_b1

with h2:
    st.markdown("### 🟡 BRAND 2 (Challenger)")
    brand2_name = st.text_input(
        "Brand Name",
        placeholder="e.g., Dr. Bombay, Häagen-Dazs...",
        key="b2",
        help=tip("brand_name"),
    )
    brand2_format = st.selectbox(
        "Product Format", 
        options=["Pint / Tub (multi-serve)", "Bag / Pouch (multi-serve)", "Box / Carton (multi-serve)", "Single-Serve Bar / Novelty", "Can / Bottle (single)"],
        index=0, 
        key="b2_format",
        help=tip("format"),
    )
    st.caption(f"**Selected: {brand2_format}**")
    st.session_state.brand2_archetype = FORMAT_TO_ARCHETYPE[brand2_format]
    st.selectbox(
        "Category",
        options=CATEGORY_OPTIONS,
        index=0,
        key="brand2_category",
    )
    
    arch2 = ARCHETYPES[FORMAT_TO_ARCHETYPE[brand2_format]]
    
    if brand2_name:
        brand_key2 = brand2_name.lower().strip()
        known_scores2 = KNOWN_BRANDS.get(brand_key2)
        
        if known_scores2:
            st.session_state.ai_scores_b2 = known_scores2
            st.success(f"✓ **{brand2_name}** — F={known_scores2['F']}, C={known_scores2['C']}")
            with st.expander("Analysis", expanded=False):
                st.write(known_scores2.get('reasoning', ''))
        elif AI_ENABLED and (st.session_state.ai_scores_b2 is None or st.session_state.last_brand2 != brand2_name):
            st.session_state.last_brand2 = brand2_name
            with st.spinner(f"Analyzing {brand2_name}..."):
                try:
                    ai_result = analyze_brand_with_ai(brand2_name)
                    if ai_result and "error" not in ai_result:
                        st.session_state.ai_scores_b2 = ai_result
                        st.success(f"✓ **{brand2_name}** — F={ai_result.get('F')}, C={ai_result.get('C')}")
                    else:
                        st.warning(f"⚠️ **{brand2_name}** not recognized — set scores manually")
                except Exception as e:
                    st.warning(f"⚠️ Set scores manually for {brand2_name}")
                    print(f"[AI ERROR B2]: {e}")
        elif st.session_state.ai_scores_b2:
            ai2 = st.session_state.ai_scores_b2
            st.success(f"✓ **{brand2_name}** — F={ai2.get('F')}, C={ai2.get('C')}")
        else:
            st.warning(f"⚠️ **{brand2_name}** — set scores manually below")
        
        ai2 = st.session_state.ai_scores_b2
        
        # === ALL 6 SCORE CONTROLS ===
        # Store original AI scores for comparison
        orig_M2 = ai2.get('M', 4) if ai2 else 4
        orig_E2 = ai2.get('E', 4) if ai2 else 4
        orig_F2 = ai2.get('F', 3) if ai2 else 3
        orig_B2 = ai2.get('B', arch2.B) if ai2 else arch2.B
        orig_K2 = ai2.get('K', arch2.K) if ai2 else arch2.K
        orig_C2 = ai2.get('C', 3) if ai2 else 3
        
        st.markdown("**ALL 6 METRICS** (adjust as needed)")
        
        # Row 1: M, E, F
        r1_c1, r1_c2, r1_c3 = st.columns(3)
        with r1_c1:
            b2_manual_M = st.number_input("M", min_value=1, max_value=5, value=orig_M2, key="b2_M", help=tip("M"))
            st.caption("Mouthfeel")
        with r1_c2:
            b2_manual_E = st.number_input("E", min_value=1, max_value=5, value=orig_E2, key="b2_E", help=tip("E"))
            st.caption("Emotion")
        with r1_c3:
            b2_manual_F = st.number_input("F ⭐", min_value=1, max_value=5, value=orig_F2, key="b2_F", help=tip("F"))
            st.caption("Familiarity")
        
        # Row 2: B, K, C
        r2_c1, r2_c2, r2_c3 = st.columns(3)
        with r2_c1:
            b2_manual_B = st.number_input("B", min_value=1, max_value=5, value=orig_B2, key="b2_B", help=tip("B"))
            st.caption("Bites")
        with r2_c2:
            b2_manual_K = st.number_input("K", min_value=1, max_value=5, value=orig_K2, key="b2_K", help=tip("K"))
            st.caption("Kinetic")
        with r2_c3:
            b2_manual_C = st.number_input("C ⭐", min_value=1, max_value=5, value=orig_C2, key="b2_C", help=tip("C"))
            st.caption("Cognitive")
        
        # Check for overrides - each needs individual justification
        b2_justifications = {}
        b2_missing_justifications = []
        
        if b2_manual_M != orig_M2:
            st.warning(f"⚠️ **M Override**: {orig_M2}→{b2_manual_M}")
            b2_justifications['M'] = st.text_input(f"Why M={b2_manual_M}?", key="b2_just_M", placeholder="Explain Mouthfeel change...", help=tip("override_rationale"))
            if len(b2_justifications['M'].strip()) < 10: b2_missing_justifications.append('M')
        
        if b2_manual_E != orig_E2:
            st.warning(f"⚠️ **E Override**: {orig_E2}→{b2_manual_E}")
            b2_justifications['E'] = st.text_input(f"Why E={b2_manual_E}?", key="b2_just_E", placeholder="Explain Emotion change...", help=tip("override_rationale"))
            if len(b2_justifications['E'].strip()) < 10: b2_missing_justifications.append('E')
        
        if b2_manual_F != orig_F2:
            st.warning(f"⚠️ **F Override**: {orig_F2}→{b2_manual_F}")
            b2_justifications['F'] = st.text_input(f"Why F={b2_manual_F}?", key="b2_just_F", placeholder="Explain Familiarity change...", help=tip("override_rationale"))
            if len(b2_justifications['F'].strip()) < 10: b2_missing_justifications.append('F')
        
        if b2_manual_B != orig_B2:
            st.warning(f"⚠️ **B Override**: {orig_B2}→{b2_manual_B}")
            b2_justifications['B'] = st.text_input(f"Why B={b2_manual_B}?", key="b2_just_B", placeholder="Explain Bites change...", help=tip("override_rationale"))
            if len(b2_justifications['B'].strip()) < 10: b2_missing_justifications.append('B')
        
        if b2_manual_K != orig_K2:
            st.warning(f"⚠️ **K Override**: {orig_K2}→{b2_manual_K}")
            b2_justifications['K'] = st.text_input(f"Why K={b2_manual_K}?", key="b2_just_K", placeholder="Explain Kinetic change...", help=tip("override_rationale"))
            if len(b2_justifications['K'].strip()) < 10: b2_missing_justifications.append('K')
        
        if b2_manual_C != orig_C2:
            st.warning(f"⚠️ **C Override**: {orig_C2}→{b2_manual_C}")
            b2_justifications['C'] = st.text_input(f"Why C={b2_manual_C}?", key="b2_just_C", placeholder="Explain Cognitive change...", help=tip("override_rationale"))
            if len(b2_justifications['C'].strip()) < 10: b2_missing_justifications.append('C')
        
        if b2_missing_justifications:
            st.error(f"❌ Justify: {', '.join(b2_missing_justifications)} (min 10 chars each)")
            st.session_state['b2_locked'] = True
        elif b2_justifications:
            st.success("✓ All overrides justified")
            st.session_state['b2_locked'] = False
        else:
            st.session_state['b2_locked'] = False
        
        b2_justification = " | ".join([f"{k}: {v}" for k, v in b2_justifications.items() if v])
        
        # Store final values
        final_scores_b2 = {
            'M': b2_manual_M, 'E': b2_manual_E, 'F': b2_manual_F,
            'B': b2_manual_B, 'K': b2_manual_K, 'C': b2_manual_C,
            'archetype': FORMAT_TO_ARCHETYPE[brand2_format],
            'reasoning': ai2.get('reasoning', 'Manual scoring') if ai2 else 'Manual scoring',
            'override_justification': b2_justification if b2_justifications else ''
        }
        st.session_state['final_b2'] = final_scores_b2

st.divider()

# ═══════════════════════════════════════════════════════════════════════════════════════
# THE STRUCTURAL COMPARISON
# ═══════════════════════════════════════════════════════════════════════════════════════

st.markdown("## ⚔️ Structural Comparison")

# Get final scores from session state (user-adjusted values)
final1 = st.session_state.get('final_b1') or {}
final2 = st.session_state.get('final_b2') or {}
arch1 = ARCHETYPES[st.session_state.brand1_archetype or "bulk"]
arch2 = ARCHETYPES[st.session_state.brand2_archetype or "bulk"]

# Brand 1 final scores - use final values if available, else archetype defaults
b1_m = final1.get('M') if final1.get('M') else arch1.M
b1_e = final1.get('E') if final1.get('E') else arch1.E
b1_f = final1.get('F') if final1.get('F') else arch1.F
b1_b = final1.get('B') if final1.get('B') else arch1.B
b1_k = final1.get('K') if final1.get('K') else arch1.K
b1_c = final1.get('C') if final1.get('C') else arch1.C

# Brand 2 final scores - use final values if available, else archetype defaults
b2_m = final2.get('M') if final2.get('M') else arch2.M
b2_e = final2.get('E') if final2.get('E') else arch2.E
b2_f = final2.get('F') if final2.get('F') else arch2.F
b2_b = final2.get('B') if final2.get('B') else arch2.B
b2_k = final2.get('K') if final2.get('K') else arch2.K
b2_c = final2.get('C') if final2.get('C') else arch2.C

# ═══════════════════════════════════════════════════════════════════════════════════════
# APPLY BEHAVIORAL PRIORS (if enabled)
# ═══════════════════════════════════════════════════════════════════════════════════════

if PRIORS_AVAILABLE and st.session_state.get('use_priors', True):
    # Build context for Brand 1
    context_b1 = {
        "cohort": st.session_state.get('cohort', 'mixed'),
        "occasion": st.session_state.get('occasion', 'evening'),
        "macro_stress": st.session_state.get('macro_stress', True),
        "category": st.session_state.get('brand1_category', DEFAULT_CATEGORY),
        "promo_frequency": st.session_state.get('promo_frequency_b1', 0.0),
        "promo_depth": st.session_state.get('promo_depth_b1', 0.0),
        "ups_pw_13": st.session_state.get('ups_pw_13_b1'),
        "ups_pw_26": st.session_state.get('ups_pw_26_b1'),
        "format": brand1_format if brand1_format else "pint",
        "use_priors": True,
    }
    
    # Build context for Brand 2
    context_b2 = {
        "cohort": st.session_state.get('cohort', 'mixed'),
        "occasion": st.session_state.get('occasion', 'evening'),
        "macro_stress": st.session_state.get('macro_stress', True),
        "category": st.session_state.get('brand2_category', DEFAULT_CATEGORY),
        "promo_frequency": st.session_state.get('promo_frequency_b2', 0.0),
        "promo_depth": st.session_state.get('promo_depth_b2', 0.0),
        "ups_pw_13": st.session_state.get('ups_pw_13_b2'),
        "ups_pw_26": st.session_state.get('ups_pw_26_b2'),
        "format": brand2_format if brand2_format else "pint",
        "use_priors": True,
    }
    
    # Apply priors to Brand 1
    raw_b1 = {"M": b1_m, "E": b1_e, "F": b1_f, "B": b1_b, "K": b1_k, "C": b1_c}
    adjusted_b1, adjustments_b1 = apply_priors_dict(raw_b1, context_b1, hard_data=HARD_DATA)
    b1_m = adjusted_b1["M"]
    b1_e = adjusted_b1["E"]
    b1_f = adjusted_b1["F"]
    b1_b = adjusted_b1["B"]
    b1_k = adjusted_b1["K"]
    b1_c = adjusted_b1["C"]
    st.session_state.adjustments_log_b1 = adjustments_b1
    
    # Apply priors to Brand 2
    raw_b2 = {"M": b2_m, "E": b2_e, "F": b2_f, "B": b2_b, "K": b2_k, "C": b2_c}
    adjusted_b2, adjustments_b2 = apply_priors_dict(raw_b2, context_b2, hard_data=HARD_DATA)
    b2_m = adjusted_b2["M"]
    b2_e = adjusted_b2["E"]
    b2_f = adjusted_b2["F"]
    b2_b = adjusted_b2["B"]
    b2_k = adjusted_b2["K"]
    b2_c = adjusted_b2["C"]
    st.session_state.adjustments_log_b2 = adjustments_b2
    
    print(f"[PRIORS] B1 adjustments: {len(adjustments_b1)} applied")
    print(f"[PRIORS] B2 adjustments: {len(adjustments_b2)} applied")
else:
    st.session_state.adjustments_log_b1 = ["Priors disabled. Using raw inputs only."]
    st.session_state.adjustments_log_b2 = ["Priors disabled. Using raw inputs only."]

# Debug output for final scores (after priors)
print(f"[DEBUG] Final B1 scores (after priors): M={b1_m}, E={b1_e}, F={b1_f}, B={b1_b}, K={b1_k}, C={b1_c}")
print(f"[DEBUG] Final B2 scores (after priors): M={b2_m}, E={b2_e}, F={b2_f}, B={b2_b}, K={b2_k}, C={b2_c}")

# Calculate S-Scores - check for locked state from overrides
b1_locked = st.session_state.get('b1_locked', False)
b2_locked = st.session_state.get('b2_locked', False)

s1 = calculate_s_score(b1_m, b1_e, b1_f, b1_b, b1_k, b1_c) if (brand1_name and not b1_locked) else None
s2 = calculate_s_score(b2_m, b2_e, b2_f, b2_b, b2_k, b2_c) if (brand2_name and not b2_locked) else None

print(f"[DEBUG] S1={s1}, S2={s2}")

# Side-by-side results
col1, col2 = st.columns(2)

with col1:
    st.markdown(f"### {brand1_name or 'Brand 1'}")
    if st.session_state.brand1_archetype == "unitized":
        st.success(f"🟢 {arch1.short_name}")
    elif st.session_state.brand1_archetype == "bulk":
        st.error(f"🔴 {arch1.short_name}")
    else:
        st.info(f"🔵 {arch1.short_name}")
    
    if brand1_name:
        st.markdown(f"""
        | Var | Score | Meaning |
        |-----|-------|---------|
        | **M** | {b1_m} | Mouthfeel |
        | **E** | {b1_e} | Emotion |
        | **F** | {b1_f} | Familiarity |
        | **B** | {b1_b} | Bites |
        | **K** | {b1_k} | Kinetic |
        | **C** | {b1_c} | Cognitive |
        """)
        num1 = b1_m * b1_e * b1_f
        den1 = b1_b * b1_k * b1_c
        st.caption(f"Numerator: {num1} | Denominator: {den1}")
        if b1_locked:
            st.error("🔒 **LOCKED**: Provide override justification above")
        elif s1:
            st.metric("S-Score™", f"{s1:.2f}")
    else:
        st.info("Enter brand name above")

with col2:
    st.markdown(f"### {brand2_name or 'Brand 2'}")
    if st.session_state.brand2_archetype == "unitized":
        st.success(f"🟢 {arch2.short_name}")
    elif st.session_state.brand2_archetype == "bulk":
        st.error(f"🔴 {arch2.short_name}")
    else:
        st.info(f"🔵 {arch2.short_name}")
    
    if brand2_name:
        st.markdown(f"""
        | Var | Score | Meaning |
        |-----|-------|---------|
        | **M** | {b2_m} | Mouthfeel |
        | **E** | {b2_e} | Emotion |
        | **F** | {b2_f} | Familiarity |
        | **B** | {b2_b} | Bites |
        | **K** | {b2_k} | Kinetic |
        | **C** | {b2_c} | Cognitive |
        """)
        num2 = b2_m * b2_e * b2_f
        den2 = b2_b * b2_k * b2_c
        st.caption(f"Numerator: {num2} | Denominator: {den2}")
        if b2_locked:
            st.error("🔒 **LOCKED**: Provide override justification above")
        elif s2:
            delta = s2 - s1 if s1 else None
            st.metric("S-Score™", f"{s2:.2f}", delta=f"{delta:+.2f}" if delta else None)
    else:
        st.info("Enter brand name above")

# Get justifications from session state for memo generation
final1 = st.session_state.get('final_b1') or {}
final2 = st.session_state.get('final_b2') or {}
rat_b1 = final1.get('override_justification', '')
rat_b2 = final2.get('override_justification', '')
b1_amendments = bool(rat_b1)
b2_amendments = bool(rat_b2)


# ═══════════════════════════════════════════════════════════════════════════════════════
# BEHAVIORAL SIGNATURE - USING NATIVE STREAMLIT
# ═══════════════════════════════════════════════════════════════════════════════════════

if s1 is not None and s2 is not None:
    st.divider()
    st.markdown("## 📊 Behavioral Signature Comparison")
    
    n1 = brand1_name or "Brand 1"
    n2 = brand2_name or "Brand 2"
    
    # Create comparison dataframe - ALL STRINGS to avoid pyarrow issues
    data = {
        "Variable": ["M · Mouthfeel", "E · Emotion", "F · Familiarity", "B · Bites", "K · Kinetic", "C · Cognitive", "S-SCORE™"],
        n1: [str(b1_m), str(b1_e), str(b1_f), str(b1_b), str(b1_k), str(b1_c), f"{s1:.2f}"],
        n2: [str(b2_m), str(b2_e), str(b2_f), str(b2_b), str(b2_k), str(b2_c), f"{s2:.2f}"],
        "Delta": [
            f"{b2_m - b1_m:+d}", f"{b2_e - b1_e:+d}", f"{b2_f - b1_f:+d}",
            f"{b2_b - b1_b:+d}", f"{b2_k - b1_k:+d}", f"{b2_c - b1_c:+d}",
            f"{s2 - s1:+.2f}"
        ]
    }
    
    df = pd.DataFrame(data)
    
    # Display as table
    st.table(df)
    
    # ═══════════════════════════════════════════════════════════════════════════════════
    # VISUAL BAR CHART COMPARISON
    # ═══════════════════════════════════════════════════════════════════════════════════
    st.markdown("### 📈 Visual Score Comparison")
    
    chart_data = pd.DataFrame({
        "Variable": ["M", "E", "F", "B", "K", "C"],
        n1: [b1_m, b1_e, b1_f, b1_b, b1_k, b1_c],
        n2: [b2_m, b2_e, b2_f, b2_b, b2_k, b2_c]
    })
    
    # Custom chart with dark theme colors
    st.bar_chart(
        chart_data.set_index("Variable"), 
        height=350, 
        use_container_width=True,
        color=["#10b981", "#ec4899"]  # Mint green and Pink for brands
    )
    
    st.caption("📈 **Numerator (M, E, F)**: Higher = more value | 📉 **Denominator (B, K, C)**: Lower = less friction")
    
    # Numerator vs Denominator metrics
    st.markdown("### Numerator vs Denominator")
    
    bc1, bc2 = st.columns(2)
    with bc1:
        st.markdown(f"**{n1}**")
        num1 = b1_m * b1_e * b1_f
        den1 = b1_b * b1_k * b1_c
        st.metric("Numerator (M×E×F)", num1)
        st.metric("Denominator (B×K×C)", den1)
    
    with bc2:
        st.markdown(f"**{n2}**")
        num2 = b2_m * b2_e * b2_f
        den2 = b2_b * b2_k * b2_c
        st.metric("Numerator (M×E×F)", num2, delta=num2-num1)
        st.metric("Denominator (B×K×C)", den2, delta=den2-den1, delta_color="inverse")

# ═══════════════════════════════════════════════════════════════════════════════════════
# INVESTOR MEMO
# ═══════════════════════════════════════════════════════════════════════════════════════

st.divider()
st.markdown("## 📋 Investor Memo: Behavioral Structural Audit")

if s1 is not None and s2 is not None:
    delta = s2 - s1
    ratio = max(s1, s2) / min(s1, s2) if min(s1, s2) > 0 else float('inf')
    winner = brand1_name if s1 > s2 else brand2_name
    
    # Big metrics - Glassmorphism style
    st.markdown(f"""
    <div style="background: rgba(91, 33, 182, 0.35); backdrop-filter: blur(12px); border: 1px solid rgba(255, 255, 255, 0.15); padding: 32px; border-radius: 20px; margin: 24px 0;">
        <div style="display: flex; justify-content: space-around; text-align: center;">
            <div>
                <div style="font-size: 0.9rem; color: #c4b5fd !important; text-transform: uppercase; letter-spacing: 0.15em; margin-bottom: 8px;">
                    {brand1_name or 'Brand 1'} S-Score™
                </div>
                <div style="font-family: 'JetBrains Mono', monospace; font-size: 3.5rem; font-weight: 800; color: #10b981 !important;">
                    {s1:.2f}
                </div>
            </div>
            <div>
                <div style="font-size: 0.9rem; color: #c4b5fd !important; text-transform: uppercase; letter-spacing: 0.15em; margin-bottom: 8px;">
                    {brand2_name or 'Brand 2'} S-Score™
                </div>
                <div style="font-family: 'JetBrains Mono', monospace; font-size: 3.5rem; font-weight: 800; color: #ec4899 !important;">
                    {s2:.2f}
                </div>
            </div>
            <div>
                <div style="font-size: 0.9rem; color: #c4b5fd !important; text-transform: uppercase; letter-spacing: 0.15em; margin-bottom: 8px;">
                    Behavioral Advantage
                </div>
                <div style="font-family: 'JetBrains Mono', monospace; font-size: 3.5rem; font-weight: 800; color: #FFFFFF !important;">
                    {ratio:.1f}x
                </div>
                <div style="font-size: 0.9rem; color: #a78bfa !important; margin-top: 4px;">{winner} leads</div>
            </div>
        </div>
    </div>
    """, unsafe_allow_html=True)
    
    n1 = brand1_name or "Brand 1"
    n2 = brand2_name or "Brand 2"
    
    # Initialize memo storage in session state
    if 'last_memo' not in st.session_state:
        st.session_state.last_memo = None
    
    scores1 = {'M': b1_m, 'E': b1_e, 'F': b1_f, 'B': b1_b, 'K': b1_k, 'C': b1_c}
    scores2 = {'M': b2_m, 'E': b2_e, 'F': b2_f, 'B': b2_b, 'K': b2_k, 'C': b2_c}
    
    # Generate memo buttons
    col_ai, col_rule = st.columns(2)
    
    with col_ai:
        if AI_ENABLED:
            if st.button("📊 Elbow Interference Investor Report", type="primary", use_container_width=True):
                with st.spinner("Generating analysis... (10-15 seconds)"):
                    try:
                        print(f"[MEMO] Attempting AI generation for {n1} vs {n2}")
                        ai_thesis = generate_strategic_synthesis(
                            n1, n2, s1, s2, scores1, scores2,
                            rat_b1 if b1_amendments else "",
                            rat_b2 if b2_amendments else ""
                        )
                        if ai_thesis:
                            st.session_state.last_memo = ai_thesis
                            print(f"[MEMO] ✓ AI memo generated successfully")
                        else:
                            print(f"[MEMO] ✗ AI returned None, using fallback")
                            st.session_state.last_memo = generate_rule_based_memo(n1, n2, s1, s2, scores1, scores2)
                            st.warning("Using standard analysis")
                    except Exception as e:
                        print(f"[MEMO] ✗ Exception: {e}")
                        st.session_state.last_memo = generate_rule_based_memo(n1, n2, s1, s2, scores1, scores2)
                        st.warning("Using standard analysis")
        else:
            st.info("Use Quick Analysis")
    
    with col_rule:
        if st.button("📋 Quick Analysis", use_container_width=True):
            st.session_state.last_memo = generate_rule_based_memo(n1, n2, s1, s2, scores1, scores2)
    
    # Display saved memo if exists
    if st.session_state.last_memo:
        st.markdown("---")
        st.markdown("### 📋 INVESTOR MEMO: BEHAVIORAL STRUCTURAL AUDIT")
        st.markdown(st.session_state.last_memo)
    
    st.divider()
    
    # Quick Summary Cards
    st.markdown("### Quick Structural Summary")
    
    sum1, sum2 = st.columns(2)
    
    with sum1:
        den1_avg = (b1_b + b1_k + b1_c) / 3
        if den1_avg <= 1.5:
            st.success(f"**{n1}**: Denominator Collapse ✓")
            st.caption("Structural Persistence — Format does the work")
        elif den1_avg <= 2.5:
            st.warning(f"**{n1}**: Partial Compression")
            st.caption("Mixed structure — Some marketing needed")
        else:
            st.error(f"**{n1}**: High Interference")
            st.caption("Purchased Velocity — Spend-dependent")
    
    with sum2:
        den2_avg = (b2_b + b2_k + b2_c) / 3
        if den2_avg <= 1.5:
            st.success(f"**{n2}**: Denominator Collapse ✓")
            st.caption("Structural Persistence — Format does the work")
        elif den2_avg <= 2.5:
            st.warning(f"**{n2}**: Partial Compression")
            st.caption("Mixed structure — Some marketing needed")
        else:
            st.error(f"**{n2}**: High Interference")
            st.caption("Purchased Velocity — Spend-dependent")
    
    # Analyst Amendment Notes
    if (b1_amendments and rat_b1.strip()) or (b2_amendments and rat_b2.strip()):
        st.divider()
        st.markdown("### Analyst Amendments")
        if b1_amendments and rat_b1.strip():
            st.info(f"**{n1}:** {rat_b1}")
        if b2_amendments and rat_b2.strip():
            st.info(f"**{n2}:** {rat_b2}")
    
    # ═══════════════════════════════════════════════════════════════════════════════════════
    # ASSUMPTIONS & ADJUSTMENTS (Priors Log)
    # ═══════════════════════════════════════════════════════════════════════════════════════
    
    if PRIORS_AVAILABLE and st.session_state.get('use_priors', True):
        st.divider()
        st.markdown("### ⚙️ Assumptions & Adjustments")
        
        with st.expander("View Prior Adjustments Applied", expanded=False):
            adj_col1, adj_col2 = st.columns(2)
            
            with adj_col1:
                st.markdown(f"**{n1} Adjustments:**")
                for adj in st.session_state.get('adjustments_log_b1', []):
                    st.caption(f"• {adj}")
            
            with adj_col2:
                st.markdown(f"**{n2} Adjustments:**")
                for adj in st.session_state.get('adjustments_log_b2', []):
                    st.caption(f"• {adj}")
            
            st.caption("---")
            st.caption(f"Context: cohort={st.session_state.get('cohort', 'mixed')}, "
                      f"occasion={st.session_state.get('occasion', 'evening')}, "
                      f"macro_stress={'Active' if st.session_state.get('macro_stress', True) else 'Inactive'}")
            
            # External Support (Hard Data)
            if HARD_DATA_AVAILABLE:
                st.markdown("---")
                st.markdown("**External Support (Hard Data)**")
                
                ext_col1, ext_col2 = st.columns(2)
                
                # Get category from context or default
                cat_b1 = st.session_state.get('brand1_category', DEFAULT_CATEGORY)
                cat_b2 = st.session_state.get('brand2_category', DEFAULT_CATEGORY)
                
                with ext_col1:
                    st.caption(f"**{n1}** ({cat_b1}):")
                    snippets_b1 = get_external_support_snippets(cat_b1)
                    for snippet in snippets_b1[:3]:  # Show top 3
                        st.caption(f"• {snippet}")
                
                with ext_col2:
                    st.caption(f"**{n2}** ({cat_b2}):")
                    snippets_b2 = get_external_support_snippets(cat_b2)
                    for snippet in snippets_b2[:3]:  # Show top 3
                        st.caption(f"• {snippet}")
                
                st.caption("---")
                st.caption("*Hard data is used as guardrails, not as score input.*")
    elif not st.session_state.get('use_priors', True):
        st.divider()
        st.caption("⚠️ **Priors disabled** — Raw inputs used without adjustment.")

else:
    st.warning("**Analysis Locked** — Provide rationale for all amendments to unlock Investor Report.")

# ═══════════════════════════════════════════════════════════════════════════════════════
# FOOTER
# ═══════════════════════════════════════════════════════════════════════════════════════

st.divider()
st.markdown("""
<div style="text-align: center; padding: 48px 0; border-top: 4px solid #10B981;">
    <div style="font-size: 1.3rem; font-weight: 900; color: #DC2626; letter-spacing: 0.25em; margin-bottom: 16px;">⬥ CONFIDENTIAL ⬥</div>
    <div style="font-size: 1.4rem; font-weight: 700; color: #475569;">PROPRIETARY METHODOLOGY OF RUSSELL BARNETT</div>
    <div style="font-size: 1.2rem; color: #64748B; margin-top: 8px;">THE ELBOW INTERFERENCE THEORY™ · ALL RIGHTS RESERVED © 2026</div>
</div>
""", unsafe_allow_html=True)
